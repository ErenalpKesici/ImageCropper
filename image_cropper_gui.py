# ImageCropper with Document to PowerPoint Converter
# Enhanced version with improved PowerPoint formatting

import time
import cv2
import numpy as np
import pytesseract
from pytesseract import Output
import re
import os
import shutil
from concurrent.futures import ThreadPoolExecutor
from functools import partial
import tkinter as tk
from tkinter import filedialog, ttk, messagebox, colorchooser
from PIL import Image, ImageTk
import threading
import multiprocessing
import requests
import io
import urllib.parse
import random
import json

# Initialize optional imports to None
Presentation = None
Inches = None
Pt = None
RGBColor = None
PP_ALIGN = None
docx = None
PyPDF2 = None
fitz = None  # PyMuPDF for better PDF reading
PPTX_AVAILABLE = False

# New imports for document to PowerPoint conversion
try:
    from pptx import Presentation as _Presentation
    from pptx.util import Inches as _Inches, Pt as _Pt
    from pptx.dml.color import RGBColor as _RGBColor
    from pptx.enum.text import PP_ALIGN as _PP_ALIGN
    import docx as _docx
    import PyPDF2 as _PyPDF2

    # Assign to module-level variables
    Presentation = _Presentation
    Inches = _Inches
    Pt = _Pt
    RGBColor = _RGBColor
    PP_ALIGN = _PP_ALIGN
    docx = _docx
    PyPDF2 = _PyPDF2
    PPTX_AVAILABLE = True
except ImportError:
    # PPTX_AVAILABLE remains False, and the variables remain None
    # The existing UI and logic checks for PPTX_AVAILABLE will handle this.
    pass

# Try to import PyMuPDF for better PDF text extraction
try:
    import fitz
    PYMUPDF_AVAILABLE = True
    print("PyMuPDF available for high-quality PDF text extraction")
except ImportError:
    PYMUPDF_AVAILABLE = False
    fitz = None
    print("PyMuPDF not available, using fallback PDF readers")

# Constants
DEFAULT_LINE_SPACING = 1.5  # Line spacing for PowerPoint slides

def contains_question_indicators(text):
    """Check if the text contains question indicators like a), A), 1), etc."""
    pattern = r'(?:^|\s)([a-zA-Z0-9])[.)]'
    matches = re.findall(pattern, text)
    unique_indicators = set(matches)
    return len(unique_indicators) >= 2

class SettingsManager:
    """Manages application settings and preferences"""
    
    def __init__(self):
        self.config_file = 'app_settings.json'
        self.default_settings = {
            'language': 'tr',
            'output_directory': os.path.join(os.getcwd(), "cropped"),
            'last_input_directory': os.getcwd(),
            'lines_per_slide': 7,
            'rows_per_page': 10,
            'detect_questions': True,            'selected_operation': 'splitter',
            'line_spacing': 1.5,
            'window_geometry': '800x900',
            'window_position': None,
            'recent_files': [],
            'recent_folders': [],
            # PowerPoint color settings (hex format)
            'title_text_color': '#000000',
            'background_color': '#ffffff', 
            'accent_color': '#1f497d'
        }
        self.settings = self.default_settings.copy()
        self.load_settings()
    
    def load_settings(self):
        """Load settings from config file"""
        try:
            if os.path.exists(self.config_file):
                with open(self.config_file, 'r', encoding='utf-8') as f:
                    saved_settings = json.load(f)
                    # Update settings with saved values, keeping defaults for missing keys
                    self.settings.update(saved_settings)
        except Exception as e:
            print(f"Error loading settings: {e}")
            self.settings = self.default_settings.copy()
    
    def save_settings(self):
        """Save current settings to config file"""
        try:
            with open(self.config_file, 'w', encoding='utf-8') as f:
                json.dump(self.settings, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"Error saving settings: {e}")
    
    def get(self, key, default=None):
        """Get a setting value"""
        return self.settings.get(key, default)
    
    def set(self, key, value):
        """Set a setting value and auto-save"""
        self.settings[key] = value
        self.save_settings()
    
    def add_recent_file(self, file_path):
        """Add a file to recent files list"""
        recent_files = self.settings.get('recent_files', [])
        if file_path in recent_files:
            recent_files.remove(file_path)
        recent_files.insert(0, file_path)
        # Keep only last 10 files
        self.settings['recent_files'] = recent_files[:10]
        self.save_settings()
    
    def add_recent_folder(self, folder_path):
        """Add a folder to recent folders list"""
        recent_folders = self.settings.get('recent_folders', [])
        if folder_path in recent_folders:
            recent_folders.remove(folder_path)
        recent_folders.insert(0, folder_path)
        # Keep only last 5 folders
        self.settings['recent_folders'] = recent_folders[:5]
        self.save_settings()

class LanguageManager:
    """Manages application language and translations"""
    
    def __init__(self, settings_manager):
        self.settings_manager = settings_manager
        self.current_language = settings_manager.get('language', 'tr')
        self.config_file = 'language_config.json'  # Keep for backward compatibility
        
        # Translation dictionaries
        self.translations = {
            'tr': {
                # Window and main labels
                'window_title': 'Resim Kırpıcı ve Belge Dönüştürücü',
                'input_section': 'Girdi',
                'output_section': 'Çıktı',
                'operation_section': 'İşlem',
                'parameters_section': 'Parametreler',
                'preview_section': 'Önizleme',
                'language_section': 'Dil',
                
                # File selection
                'no_files_selected': 'Dosya seçilmedi',
                'images_selected': 'resim seçildi',
                'documents_selected': 'belge seçildi',
                'images_from_folder': 'klasörden resim',
                'select_folder': 'Klasör Seç',
                'select_images': 'Resim Seç',
                'select_documents': 'Belge Seç',
                'select_output_folder': 'Klasör Seç',
                
                # Operations
                'split_images': 'Resimleri Böl',
                'crop_titles': 'Başlıkları Kırp',
                'remove_blank_images': 'Boş Resimleri Kaldır',
                'convert_doc_to_pptx': 'Belgeyi PowerPoint\'e Dönüştür',
                  # Parameters
                'lines_per_slide': 'Slayt başına satır (belge dönüşümü için):',
                'lines_recommendation': 'Önerilen: Taşma olmaması için 5-8 satır, büyük fontlar',
                'line_spacing': 'Satır aralığı:',
                'line_spacing_recommendation': 'Önerilen: 1.2-2.0 arası (varsayılan: 1.5)',
                'title_text_color': 'Başlık metin rengi:',
                'background_color': 'Arka plan rengi:',
                'accent_color': 'Vurgu rengi:',
                'choose_color': 'Renk Seçin',                'rows_per_page': 'Sayfa başına satır (resim bölme için):',
                'max_splits': 'Maksimum bölme sayısı:',
                'max_splits_help': 'Resmi kaç parçaya böleceğinizi belirler (1-5 arası)',
                'detect_questions': 'Soruları algıla (bölme)',
                
                # Color settings
                'color_settings': 'Renk Ayarları',
                'title_text_color': 'Başlık metni rengi:',
                'content_text_color': 'İçerik metni rengi:',
                'title_background_color': 'Başlık arka plan rengi:',
                'slide_background_color': 'Slayt arka plan rengi:',
                'accent_color': 'Vurgu rengi:',
                'table_header_color': 'Tablo başlığı rengi:',
                'choose_color': 'Renk Seç',
                'reset_colors': 'Renkleri Sıfırla',
                
                # Buttons
                'process': 'İşle',
                'cancel': 'İptal',
                'open_output_folder': 'Çıktı Klasörünü Aç',
                
                # Status messages
                'ready': 'Hazır',
                'processing': 'İşleniyor...',
                'cancelling': 'İptal ediliyor...',
                'processing_complete': 'İşlem tamamlandı!',
                'document_conversion_complete': 'Belge dönüşümü tamamlandı!',
                
                # Error messages
                'missing_libraries': 'Eksik Kütüphaneler',
                'missing_libraries_msg': 'Belge dönüşümü için python-pptx, python-docx ve PyPDF2 gereklidir.\nKurulum: pip install python-pptx python-docx PyPDF2',
                'no_input': 'Girdi Yok',
                'no_input_msg': 'Önce dosya seçin.',
                'no_images': 'Resim Yok',
                'no_images_msg': 'Seçilen klasörde resim dosyası bulunamadı.',
                'folder_not_found': 'Klasör Bulunamadı',
                'folder_not_found_msg': 'Çıktı klasörü henüz mevcut değil.',
                'error': 'Hata',
                'error_occurred': 'Bir hata oluştu: ',
                
                # Success messages
                'complete': 'Tamamlandı',
                'successfully_processed': 'başarıyla işlendi',
                'successfully_converted': 'başarıyla PowerPoint\'e dönüştürüldü',
                'found_blank_images': 'boş resim bulundu',
                'blank_images_found': 'resim içinden boş resim bulundu',
                
                # File processing
                'processed_files': 'işlenen dosya',
                'reading_file': 'dosyası okunuyor ve tablolar algılanıyor...',
                'creating_slides': 'Tablolar ve optimal boşluklarla slaytlar oluşturuluyor...',
                'created_powerpoint': 'PowerPoint oluşturuldu:',
                'slides_created': 'slayt ile',
                'including_tables': 'tablo(lar) dahil',
                'processing_error': 'İşleme hatası',
                  # Language
                'language': 'Dil:',
                'settings': 'Ayarlar',
                'turkish': 'Türkçe',
                'english': 'English'
            },
            'en': {
                # Window and main labels
                'window_title': 'Image Cropper & Document Converter',
                'input_section': 'Input',
                'output_section': 'Output',
                'operation_section': 'Operation',
                'parameters_section': 'Parameters',
                'preview_section': 'Preview',
                'language_section': 'Language',
                
                # File selection
                'no_files_selected': 'No files selected',
                'images_selected': 'image(s) selected',
                'documents_selected': 'document(s) selected',
                'images_from_folder': 'image(s) from folder',
                'select_folder': 'Select Folder',
                'select_images': 'Select Images',
                'select_documents': 'Select Documents',
                'select_output_folder': 'Select Folder',
                
                # Operations
                'split_images': 'Split Images',
                'crop_titles': 'Crop Titles',
                'remove_blank_images': 'Remove Blank Images',
                'convert_doc_to_pptx': 'Convert Document to PowerPoint',
                  # Parameters
                'lines_per_slide': 'Lines per slide (for document conversion):',
                'lines_recommendation': 'Recommended: 5-8 lines for no overflow, larger fonts',
                'line_spacing': 'Line spacing:',
                'line_spacing_recommendation': 'Recommended: 1.2-2.0 range (default: 1.5)',
                'title_text_color': 'Title text color:',
                'background_color': 'Background color:',
                'accent_color': 'Accent color:',
                'choose_color': 'Choose Color',                'rows_per_page': 'Rows per page (for image splitting):',
                'max_splits': 'Maximum number of splits:',
                'max_splits_help': 'How many parts to split the image into (1-5 range)',
                'detect_questions': 'Detect questions (splitting)',
                
                # Color settings
                'color_settings': 'Color Settings',
                'title_text_color': 'Title text color:',
                'content_text_color': 'Content text color:',
                'title_background_color': 'Title background color:',
                'slide_background_color': 'Slide background color:',
                'accent_color': 'Accent color:',
                'table_header_color': 'Table header color:',
                'choose_color': 'Choose Color',
                'reset_colors': 'Reset Colors',
                
                # Buttons
                'process': 'Process',
                'cancel': 'Cancel',
                'open_output_folder': 'Open Output Folder',
                
                # Status messages
                'ready': 'Ready',
                'processing': 'Processing...',
                'cancelling': 'Cancelling...',
                'processing_complete': 'Processing complete!',
                'document_conversion_complete': 'Document conversion complete!',
                
                # Error messages
                'missing_libraries': 'Missing Libraries',
                'missing_libraries_msg': 'Document conversion requires python-pptx, python-docx, and PyPDF2.\nInstall with: pip install python-pptx python-docx PyPDF2',
                'no_input': 'No Input',
                'no_input_msg': 'Please select files first.',
                'no_images': 'No Images',
                'no_images_msg': 'No image files found in the selected folder.',
                'folder_not_found': 'Folder Not Found',
                'folder_not_found_msg': 'Output folder does not exist yet.',
                'error': 'Error',
                'error_occurred': 'An error occurred: ',
                
                # Success messages
                'complete': 'Complete',
                'successfully_processed': 'successfully processed',
                'successfully_converted': 'successfully converted',
                'found_blank_images': 'blank images found',
                'blank_images_found': 'blank images found out of',
                
                # File processing
                'processed_files': 'processed files',
                'reading_file': 'file and detecting tables...',
                'creating_slides': 'Creating slides with tables and optimal spacing...',
                'created_powerpoint': 'Created PowerPoint:',
                'slides_created': 'with slides',
                'including_tables': 'including table(s)',
                'processing_error': 'Processing error',
                  # Language
                'language': 'Language:',
                'settings': 'Settings',
                'turkish': 'Türkçe',
                'english': 'English'
            }
        }
    
    def get_text(self, key):
        """Get translated text for the current language"""
        text = self.translations.get(self.current_language, {}).get(key, key)
        return text if text is not None else key  # Ensure we never return None
    
    def set_language(self, language_code):
        """Set the current language"""
        if language_code in self.translations:
            self.current_language = language_code
            self.settings_manager.set('language', language_code)  # Use settings manager instead
    
    def get_current_language(self):
        """Get the current language code"""
        return self.current_language
    
    def get_available_languages(self):
        """Get list of available language codes"""
        return list(self.translations.keys())
    
    def load_language_preference(self):
        """Load language preference from config file - kept for backward compatibility"""
        try:
            if os.path.exists(self.config_file):
                with open(self.config_file, 'r', encoding='utf-8') as f:
                    config = json.load(f)
                    self.current_language = config.get('language', 'tr')
        except Exception as e:
            print(f"Error loading language preference: {e}")
            self.current_language = 'tr'  # Default fallback
    
    def save_language_preference(self):
        """Save language preference to config file - kept for backward compatibility"""
        try:
            config = {'language': self.current_language}
            with open(self.config_file, 'w', encoding='utf-8') as f:
                json.dump(config, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"Error saving language preference: {e}")

def remove_empty_rows_and_columns(image, empty_threshold=100):
    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    row_sums = np.sum(gray < empty_threshold, axis=1)
    col_sums = np.sum(gray < empty_threshold, axis=0)
    
    non_empty_rows = np.where(row_sums > 0)[0]
    non_empty_cols = np.where(col_sums > 0)[0]
    
    if len(non_empty_rows) > 0 and len(non_empty_cols) > 0:
        image = image[non_empty_rows[0]:non_empty_rows[-1]+1, non_empty_cols[0]:non_empty_cols[-1]+1]
    return image

def add_outer_border(image, top_border=10, bottom_border=10, left_border=20, right_border=20):
    outer_color = image[0, 0].tolist()
    bordered_image = cv2.copyMakeBorder(
        image,
        top_border,
        bottom_border,
        left_border,
        right_border,
        cv2.BORDER_CONSTANT,
        value=outer_color
    )
    return bordered_image

def is_image_blank(image, background_threshold=99.0, min_contour_area=500, std_dev_threshold=10):
    """Check if an image is effectively blank"""
    pixels = image.reshape(-1, 3)
    pixel_values = (pixels[:, 0] << 16) | (pixels[:, 1] << 8) | pixels[:, 2]
    unique_values, counts = np.unique(pixel_values, return_counts=True)
    background_value = unique_values[np.argmax(counts)]
    
    background_color = np.array([
        (background_value >> 16) & 255,
        (background_value >> 8) & 255,
        background_value & 255
    ])
    
    color_distance = np.sqrt(np.sum((pixels - background_color)**2, axis=1))
    background_pixels = np.sum(color_distance < 10)
    background_percentage = (background_pixels / pixels.shape[0]) * 100
    
    if background_percentage >= background_threshold:
        return True
    
    std_dev = np.std(pixels, axis=0)
    if np.all(std_dev < std_dev_threshold):
        return True
    
    text = pytesseract.image_to_string(image)
    if text == '' or re.fullmatch(r'\d+', text.strip()):
        return True
    
    return False

class ImageCropperApp:
    def __init__(self, root):
        self.root = root
        
        # Initialize settings manager first
        self.settings = SettingsManager()
          # Initialize language manager with settings
        self.lang = LanguageManager(self.settings)
          # Set window properties from settings
        self.root.title(self.lang.get_text('window_title'))
        geometry = self.settings.get('window_geometry', '800x900')
        
        # Load saved window position if available
        position = self.settings.get('window_position')
        if position:
            # position should be in format "x+y" like "+100+200"
            if not position.startswith('+'):
                position = f"+{position}"
            self.root.geometry(f"{geometry}{position}")
        else:
            self.root.geometry(geometry)
        
        self.root.minsize(800, 900)
        
        # Initialize variables from settings
        self.input_files = []
        self.output_dir = self.settings.get('output_directory', os.path.join(os.getcwd(), "cropped"))
        self.last_title = ''
        self.last_top = 0
        self.processing_cancelled = False
        
        # Dictionary to store widget references for language updates
        self.widgets = {}        
        # Bind window close event to save settings
        self.root.protocol("WM_DELETE_WINDOW", self.on_closing)
        
        # Initialize language variable for menu
        self.language_var = tk.StringVar(value=self.lang.get_current_language())
        
        self.create_widgets()
        # Load saved values into widgets
        self.load_saved_values()
        
    def create_menu_bar(self):
        """Create menu bar with language selection"""
        menubar = tk.Menu(self.root)
        self.root.config(menu=menubar)
        
        # Settings menu
        settings_menu = tk.Menu(menubar, tearoff=0)
        menubar.add_cascade(label=self.lang.get_text('settings'), menu=settings_menu)
        
        # Language submenu
        language_menu = tk.Menu(settings_menu, tearoff=0)
        settings_menu.add_cascade(label=self.lang.get_text('language'), menu=language_menu)
        
        # Add language options
        language_menu.add_radiobutton(
            label="Türkçe", 
            variable=self.language_var, 
            value="tr",
            command=self.on_language_change_menu
        )
        language_menu.add_radiobutton(
            label="English", 
            variable=self.language_var, 
            value="en",
            command=self.on_language_change_menu
        )        # Store menu references for language updates
        self.widgets['menubar'] = menubar
        self.widgets['settings_menu'] = settings_menu
        self.widgets['language_menu'] = language_menu
    
    def on_language_change_menu(self):
        """Handle language change from menu"""
        new_language = self.language_var.get()
        self.lang.set_language(new_language)
        self.update_all_ui_text()
        
        # Save the language change
        self.settings.set('language', new_language)
        
        # Recreate menu to update labels properly
        self.create_menu_bar()
    
    def create_widgets(self):
        # Create menu bar
        self.create_menu_bar()
        
        # Create main frame with scrolling
        main_frame = ttk.Frame(self.root, padding="10")
        main_frame.pack(fill=tk.BOTH, expand=True)
        
        canvas = tk.Canvas(main_frame)
        scrollbar = ttk.Scrollbar(main_frame, orient="vertical", command=canvas.yview)
        scrollable_frame = ttk.Frame(canvas)
        
        scrollable_frame.bind(
            "<Configure>",
            lambda e: canvas.configure(scrollregion=canvas.bbox("all"))
        )
        
        canvas.create_window((0, 0), window=scrollable_frame, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)
        
        canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")
          # Input section
        input_frame = ttk.LabelFrame(scrollable_frame, text=self.lang.get_text('input_section'), padding="10")
        input_frame.pack(fill=tk.X, pady=5)
        self.widgets['input_frame'] = input_frame
        
        self.input_files_label = ttk.Label(input_frame, text=self.lang.get_text('no_files_selected'))
        self.input_files_label.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        self.widgets['input_files_label'] = self.input_files_label
        
        # Button frame for file selection
        select_buttons_frame = ttk.Frame(input_frame)
        select_buttons_frame.pack(side=tk.RIGHT)
        
        select_folder_button = ttk.Button(select_buttons_frame, text=self.lang.get_text('select_folder'), command=self.select_folder)
        select_folder_button.pack(side=tk.LEFT, padx=5)
        self.widgets['select_folder_button'] = select_folder_button
        
        select_images_button = ttk.Button(select_buttons_frame, text=self.lang.get_text('select_images'), command=self.select_files)
        select_images_button.pack(side=tk.LEFT, padx=5)
        self.widgets['select_images_button'] = select_images_button
        
        select_docs_button = ttk.Button(select_buttons_frame, text=self.lang.get_text('select_documents'), command=self.select_documents)
        select_docs_button.pack(side=tk.LEFT, padx=5)
        self.widgets['select_docs_button'] = select_docs_button
        
        # Output section
        output_frame = ttk.LabelFrame(scrollable_frame, text=self.lang.get_text('output_section'), padding="10")
        output_frame.pack(fill=tk.X, pady=5)
        self.widgets['output_frame'] = output_frame
        
        self.output_dir_label = ttk.Label(output_frame, text=self.output_dir)
        self.output_dir_label.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        
        output_button = ttk.Button(output_frame, text=self.lang.get_text('select_output_folder'), command=self.select_output_dir)
        output_button.pack(side=tk.RIGHT, padx=5)
        self.widgets['output_button'] = output_button
        
        # Operation section
        operation_frame = ttk.LabelFrame(scrollable_frame, text=self.lang.get_text('operation_section'), padding="10")
        operation_frame.pack(fill=tk.X, pady=5)
        self.widgets['operation_frame'] = operation_frame
        
        self.operation_var = tk.StringVar(value="splitter")
        # Add trace to handle operation changes
        self.operation_var.trace('w', self.on_operation_change)
        
        # Store radio buttons for updating
        self.operation_radios = []
        
        operations = [
            ("split_images", "splitter"), 
            ("crop_titles", "title_cropper"),
            ("remove_blank_images", "blank_remover"),
            ("convert_doc_to_pptx", "doc_to_pptx")
        ]
        
        for text_key, value in operations:
            radio = ttk.Radiobutton(operation_frame, text=self.lang.get_text(text_key), 
                                  value=value, variable=self.operation_var)
            radio.pack(anchor=tk.W)
            self.operation_radios.append((radio, text_key))
          # Parameters section
        param_frame = ttk.LabelFrame(scrollable_frame, text=self.lang.get_text('parameters_section'), padding="10")
        param_frame.pack(fill=tk.X, pady=5)
        self.widgets['param_frame'] = param_frame
        
        # Dynamic fields that will be shown/hidden based on operation
        self.dynamic_widgets = {}
        
        # Create all possible dynamic fields
        self.create_dynamic_fields(param_frame)
        
        # Initialize field visibility based on default operation
        self.on_operation_change()
        
        # Preview area
        preview_frame = ttk.LabelFrame(scrollable_frame, text=self.lang.get_text('preview_section'), padding="10")
        preview_frame.pack(fill=tk.BOTH, expand=True, pady=5)
        self.widgets['preview_frame'] = preview_frame
        
        self.preview_area = ttk.Frame(preview_frame)
        self.preview_area.pack(fill=tk.BOTH, expand=True)
          # Progress and status
        self.progress_var = tk.DoubleVar(value=0)
        progress_bar = ttk.Progressbar(scrollable_frame, variable=self.progress_var, maximum=100)
        progress_bar.pack(fill=tk.X, pady=5)
        
        self.status_var = tk.StringVar(value=self.lang.get_text('ready'))
        status_label = ttk.Label(scrollable_frame, textvariable=self.status_var)
        status_label.pack(anchor=tk.W, pady=5)
        
        # Action buttons
        button_frame = ttk.Frame(scrollable_frame)
        button_frame.pack(fill=tk.X, pady=10)
        
        process_button = ttk.Button(button_frame, text=self.lang.get_text('process'), command=self.process_files)
        process_button.pack(side=tk.LEFT, padx=5)
        self.widgets['process_button'] = process_button
        
        self.cancel_button = ttk.Button(button_frame, text=self.lang.get_text('cancel'), command=self.cancel_processing, state=tk.DISABLED)
        self.cancel_button.pack(side=tk.LEFT, padx=5)
        self.widgets['cancel_button'] = self.cancel_button
        
        open_folder_button = ttk.Button(button_frame, text=self.lang.get_text('open_output_folder'), command=self.open_output_folder)
        open_folder_button.pack(side=tk.RIGHT, padx=5)
        self.widgets['open_folder_button'] = open_folder_button
    
    def select_files(self):
        """Select image files"""
        initial_dir = self.settings.get('last_input_directory', os.getcwd())
        filetypes = [("Image files", "*.png *.jpg *.jpeg *.bmp *.tiff")]
        files = filedialog.askopenfilenames(filetypes=filetypes, initialdir=initial_dir)
        
        if files:
            self.input_files = list(files)
            text = f"{len(self.input_files)} {self.lang.get_text('images_selected')}"
            self.input_files_label.config(text=text)
            self.update_preview()
            
            # Save the directory for future use
            directory = os.path.dirname(files[0])
            self.settings.set('last_input_directory', directory)
            
            # Add files to recent files
            for file_path in files[:5]:  # Only save first 5 to avoid too many entries
                self.settings.add_recent_file(file_path)
    
    def select_documents(self):
        """Select Word or PDF documents for conversion to PowerPoint"""
        if not PPTX_AVAILABLE:
            messagebox.showerror(self.lang.get_text('missing_libraries'), 
                self.lang.get_text('missing_libraries_msg'))
            return
        
        initial_dir = self.settings.get('last_input_directory', os.getcwd())
        filetypes = [("Document files", "*.docx *.pdf"), ("Word files", "*.docx"), ("PDF files", "*.pdf")]
        files = filedialog.askopenfilenames(filetypes=filetypes, initialdir=initial_dir)
        
        if files:
            self.input_files = list(files)
            text = f"{len(self.input_files)} {self.lang.get_text('documents_selected')}"
            self.input_files_label.config(text=text)
            self.update_document_preview()
            
            # Save the directory for future use
            directory = os.path.dirname(files[0])
            self.settings.set('last_input_directory', directory)
            
            # Add files to recent files
            for file_path in files[:5]:  # Only save first 5 to avoid too many entries
                self.settings.add_recent_file(file_path)
    
    def select_folder(self):
        """Select a folder and add all images from it"""
        initial_dir = self.settings.get('last_input_directory', os.getcwd())
        folder = filedialog.askdirectory(initialdir=initial_dir)
        if folder:
            image_files = []
            for file in os.listdir(folder):
                if file.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tiff')):
                    image_files.append(os.path.join(folder, file))
                    
            if image_files:
                self.input_files = image_files
                text = f"{len(self.input_files)} {self.lang.get_text('images_from_folder')}"
                self.input_files_label.config(text=text)
                self.update_preview()
                
                # Save the directory for future use
                self.settings.set('last_input_directory', folder)
                self.settings.add_recent_folder(folder)
            else:
                messagebox.showinfo(self.lang.get_text('no_images'), 
                    self.lang.get_text('no_images_msg'))
    
    def select_output_dir(self):
        initial_dir = self.settings.get('output_directory', os.getcwd())
        directory = filedialog.askdirectory(initialdir=initial_dir)
        if directory:
            self.output_dir = directory
            self.output_dir_label.config(text=self.output_dir)
            # Save automatically
            self.settings.set('output_directory', directory)
            self.settings.add_recent_folder(directory)
    
    def update_preview(self):
        """Show thumbnails of selected images"""
        for widget in self.preview_area.winfo_children():
            widget.destroy()
            
        max_previews = min(5, len(self.input_files))
        self.thumbnail_refs = []
        
        for i in range(max_previews):
            try:
                img = Image.open(self.input_files[i])
                img.thumbnail((150, 150))
                photo = ImageTk.PhotoImage(img)
                
                self.thumbnail_refs.append(photo)
                
                frame = ttk.Frame(self.preview_area)
                frame.pack(side=tk.LEFT, padx=5)
                
                label = ttk.Label(frame, image=photo)
                label.pack()
                
                name_label = ttk.Label(frame, text=os.path.basename(self.input_files[i]), wraplength=150)
                name_label.pack()
            except Exception as e:
                print(f"Error creating thumbnail: {e}")
    def update_document_preview(self):
        """Show document names for preview"""
        for widget in self.preview_area.winfo_children():
            widget.destroy()
            
        for i, doc_path in enumerate(self.input_files[:5]):
            frame = ttk.Frame(self.preview_area)
            frame.pack(side=tk.LEFT, padx=5)
            
            name_label = ttk.Label(frame, text=os.path.basename(doc_path), wraplength=150)
            name_label.pack()
            
            ext = os.path.splitext(doc_path)[1].upper()
            type_label = ttk.Label(frame, text=f"{ext} Document", font=("Arial", 8))
            type_label.pack()
    
    def convert_document_to_pptx(self, document_path, lines_per_slide):
        """Convert Word or PDF document to PowerPoint presentation with table detection"""
        if not PPTX_AVAILABLE:
            self.status_var.set("PowerPoint libraries not available")
            return False
            
        try:
            # Get line spacing from dynamic field if available
            line_spacing = DEFAULT_LINE_SPACING
            if hasattr(self, 'line_spacing_var'):
                try:
                    line_spacing = float(self.line_spacing_var.get())
                    # Validate range
                    if line_spacing < 0.8 or line_spacing > 3.0:
                        line_spacing = DEFAULT_LINE_SPACING
                except (ValueError, AttributeError):
                    line_spacing = DEFAULT_LINE_SPACING
            
            content_items = []
            file_ext = os.path.splitext(document_path)[1].lower()
            self.status_var.set(f"{file_ext} {self.lang.get_text('reading_file')}")
            self.root.update()
            
            if file_ext == '.docx':
                content_items = self.extract_docx_content_with_tables(document_path)
            
            elif file_ext == '.pdf':
                content_items = self.extract_pdf_content_with_tables(document_path)
            if not content_items:
                self.status_var.set(f"No content found in {document_path}")
                return False
              # Create PowerPoint presentation
            if not PPTX_AVAILABLE:
                self.status_var.set("PowerPoint libraries not available")
                return False
                
            prs = Presentation()
            if Inches is not None:
                prs.slide_width = Inches(16)
                prs.slide_height = Inches(9)

            slide_layout = prs.slide_layouts[6]  # Blank layout
            
            current_slide_lines = []
            slide_count = 0
            table_count = 0
            self.status_var.set(self.lang.get_text('creating_slides'))
            self.root.update()
            
            # Process content items (text and tables)
            for item in content_items:
                if item['type'] == 'table':
                    # Create a dedicated slide for the table
                    if current_slide_lines:
                        # Finish current text slide first
                        slide_count += 1
                        self.create_optimized_slide(prs, slide_layout, current_slide_lines, line_spacing)
                        current_slide_lines = []
                    
                    # Create table slide
                    slide_count += 1
                    table_count += 1
                    self.create_table_slide(prs, slide_layout, item['content'])
                    
                elif item['type'] == 'text':                    # Add text to current slide
                    current_slide_lines.append(item['content'])
                    
                    # Create new slide when we reach the line limit
                    if len(current_slide_lines) >= lines_per_slide:
                        slide_count += 1
                        self.create_optimized_slide(prs, slide_layout, current_slide_lines, line_spacing)
                        current_slide_lines = []

            # Add remaining lines to final slide
            if current_slide_lines:
                slide_count += 1
                self.create_optimized_slide(prs, slide_layout, current_slide_lines, line_spacing)
              # Save the presentation
            base_name = os.path.splitext(os.path.basename(document_path))[0]
            output_path = os.path.join(self.output_dir, f"{base_name}.pptx")
            prs.save(output_path)
            
            status_msg = f"{self.lang.get_text('created_powerpoint')} {base_name}.pptx {self.lang.get_text('slides_created')} {slide_count}"
            if table_count > 0:
                status_msg += f" ({self.lang.get_text('including_tables')} {table_count})"
            self.status_var.set(status_msg)
            return True
        except Exception as e:
            error_msg = f"{self.lang.get_text('processing_error')} {document_path}: {str(e)}"
            self.status_var.set(error_msg)
            print(f"Error converting document {document_path}: {e}")
            return False
    
    def calculate_text_height(self, lines, font_size, line_spacing=1.5):
        """Calculate more accurate text height based on font size and line count"""
        # More precise calculation considering PowerPoint's text rendering
        line_height_pt = font_size * line_spacing
        total_height_pt = len(lines) * line_height_pt
        
        # Add space_after for each line (varies by content type)
        space_after_pt = len(lines) * 2  # Average space_after
          # Add some padding for text frame margins and safety buffer
        padding_pt = 20  # Safety margin
        
        total_height = total_height_pt + space_after_pt + padding_pt
        return total_height
    def find_optimal_font_size(self, lines, available_height_inches, text_type='content', line_spacing=1.5):
        """Find the maximum font size that fits without overflow using binary search"""
        available_height_pt = available_height_inches * 72  # Convert inches to points
        
        # Define size ranges based on text type - increased max sizes for better utilization
        min_size_arr = [('title', 22), ('subtitle', 20), ('bullet', 16), ('content', 14)]
        min_size_dict = dict(min_size_arr)
        min_size = min_size_dict[text_type]
        max_size = min_size * 3
        
        # Use binary search for more efficient and precise font sizing
        left, right = min_size, max_size
        optimal_size = min_size
        
        while left <= right:
            mid = (left + right) // 2
            estimated_height = self.calculate_text_height(lines, mid, line_spacing)
            
            if estimated_height <= available_height_pt:
                optimal_size = mid  # This size fits, try larger
                left = mid + 1
            else:
                right = mid - 1  # This size is too big, try smaller
        
        # Fine-tune with smaller increments around the optimal size
        for size in range(max(optimal_size - 2, min_size), min(optimal_size + 3, max_size + 1)):
            estimated_height = self.calculate_text_height(lines, size, line_spacing)
            if estimated_height <= available_height_pt:
                optimal_size = max(optimal_size, size)
        
        print(f"📏 Font sizing: {len(lines)} lines, {text_type} -> {optimal_size}pt (range: {min_size}-{max_size}pt)")
        return optimal_size
    
    def detect_and_fix_overflow(self, text_frame, available_height_inches):
        """Detect text overflow and automatically reduce font sizes to fit"""
        available_height_pt = available_height_inches * 72
        max_attempts = 8  # Prevent infinite loops
        attempt = 0
        
        while attempt < max_attempts:
            # Calculate actual text height from the text frame
            actual_height_pt = self.estimate_actual_text_height(text_frame)
            
            # Add safety margin (10% of available height)
            safety_margin = available_height_pt * 0.1
            
            if actual_height_pt <= (available_height_pt - safety_margin):
                if attempt > 0:
                    print(f"✅ Text fits after {attempt} adjustments: {actual_height_pt:.1f}pt <= {available_height_pt:.1f}pt")
                break
            
            print(f"⚠️  Overflow detected (attempt {attempt + 1}): {actual_height_pt:.1f}pt > {available_height_pt:.1f}pt")
            
            # Reduce all font sizes by 15% per iteration
            reduction_factor = 0.85
            adjusted = False
            
            for paragraph in text_frame.paragraphs:
                if paragraph.font.size:
                    current_size = paragraph.font.size.pt
                    new_size = max(8, int(current_size * reduction_factor))  # Minimum 8pt
                    if new_size != current_size:
                        paragraph.font.size = Pt(new_size)
                        adjusted = True
                        if attempt == 0:  # Only print details on first adjustment
                            print(f"🔽 Reduced font: {current_size:.0f}pt -> {new_size}pt")
            
            if not adjusted:
                print("⚠️  Could not reduce fonts further (minimum reached)")
                break
                
            attempt += 1
        
        if attempt >= max_attempts:
            print(f"⚠️  Warning: Max adjustment attempts reached, some overflow may remain")
    
    def estimate_actual_text_height(self, text_frame):
        """Estimate the actual rendered text height in PowerPoint"""
        total_height = 0
        line_count = 0
        
        for paragraph in text_frame.paragraphs:
            if paragraph.text.strip():
                # Get font size (default to 12pt if not set)
                font_size = 12
                if paragraph.font.size:
                    font_size = paragraph.font.size.pt
                  # Calculate line height based on font size and line spacing
                line_spacing = getattr(paragraph, 'line_spacing', 1.5)
                if isinstance(line_spacing, (int, float)):
                    line_height = font_size * line_spacing
                else:
                    line_height = font_size * 1.5  # Default line spacing
                
                # Add space after paragraph
                space_after = 0
                if paragraph.space_after:
                    space_after = paragraph.space_after.pt
                
                # Count wrapped lines (rough estimation)
                text_length = len(paragraph.text)
                estimated_chars_per_line = max(50, 100 - (font_size - 12) * 2)  # Larger fonts = fewer chars per line
                estimated_lines = max(1, text_length // estimated_chars_per_line)
                
                paragraph_height = (line_height * estimated_lines) + space_after
                total_height += paragraph_height
                line_count += estimated_lines
        
        # Add text frame margins
        margin_top = 0
        margin_bottom = 0
        if hasattr(text_frame, 'margin_top') and text_frame.margin_top:
            margin_top = text_frame.margin_top / 914400 * 72  # Convert EMU to points
        if hasattr(text_frame, 'margin_bottom') and text_frame.margin_bottom:
            margin_bottom = text_frame.margin_bottom / 914400 * 72
        total_height += margin_top + margin_bottom        
        return total_height
    
    def create_optimized_slide(self, prs, slide_layout, lines, line_spacing=1.5):
        """Create slide with dynamic font sizing for maximum content optimization"""
        slide = prs.slides.add_slide(slide_layout)
        
        # Set slide background color if custom color is specified
        bg_rgb = self.get_color_setting('background_color', '#ffffff')
        if bg_rgb != (255, 255, 255):  # Only set if not default white
            try:
                # Set slide background color
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*bg_rgb)
            except Exception as e:
                print(f"Could not set background color: {e}")
        
        # Use FULL slide dimensions for maximum space utilization
        left = Inches(0.3)      # Minimal left margin
        top = 0 # Inches(0.3)       # Minimal top margin  
        width = Inches(16)    
        height = Inches(9)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = None  # Critical: prevent auto-sizing
        text_frame.line_spacing = line_spacing
        
        # Minimal margins for maximum space
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.0)
        text_frame.margin_bottom = Inches(0.0)
        
        # Available height for text (accounting for minimal margins)
        available_height = height - Inches(0.1)  # Subtract tiny margins
          # Check if first line is a title using improved logic
        is_title = self.is_better_title(lines[0] if lines else "", 0, lines)
        
        # Calculate optimal font sizes for different content types
        title_lines = [lines[0]] if is_title and len(lines) > 0 else []
        content_lines = lines[1:] if is_title else lines
          # Find optimal font sizes
        title_font_size = 24  # Default title font size
        if title_lines:
            title_font_size = self.find_optimal_font_size([title_lines[0]], available_height * 0.25, 'title', line_spacing)
        
        base_font_size = 16  # Default base font size
        if content_lines:
            # Reserve space for title if present
            content_height = available_height * 0.75 if title_lines else available_height
            
            # Calculate base font size for all content combined
            base_font_size = self.find_optimal_font_size(content_lines, content_height, 'content', line_spacing)
        
        print(f"🎯 Slide with {len(lines)} lines - Title: {title_font_size}pt, Content: {base_font_size}pt")
        
        for i, line in enumerate(lines):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
              # Enhance text with Unicode symbols where appropriate
            is_current_line_title = (i == 0 and is_title)
            enhanced_text = self.add_unicode_symbols_to_text(line, is_current_line_title)
            p.text = enhanced_text
            p.font.name = 'Calibri'
            p.word_wrap = True
            
            # Dynamic font sizing based on content type and available space
            if i == 0 and is_title:                # Title: Dynamic sizing with emphasis
                p.font.size = Pt(max(title_font_size, 18))  # Minimum 18pt for titles
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
                title_rgb = self.get_color_setting('title_text_color', '#1f497d')
                p.font.color.rgb = RGBColor(*title_rgb)
                p.space_after = Pt(max(4, title_font_size // 4))
            elif line.startswith(('•', '-', '*', '▪', '○')) or any(line.startswith(f'{j}.') for j in range(1, 50)):
                # Bullet points: Slightly smaller than regular content
                p.font.size = Pt(max(base_font_size - 2, 12))
                p.font.bold = False
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(2)
            elif len(line) < 60 and not line.endswith(('.', '!', '?')) and i > 0:
                # Subtitles: Larger than regular content
                p.font.size = Pt(max(base_font_size + 2, 14))
                p.font.bold = True
                p.alignment = PP_ALIGN.LEFT
                p.font.color.rgb = RGBColor(68, 84, 106)
                p.space_after = Pt(3)
            else:
                # Regular content: Optimized dynamic size
                p.font.size = Pt(max(base_font_size, 12))  # Minimum 12pt for readability
                p.font.bold = False
                p.alignment = PP_ALIGN.LEFT                
                p.space_after = Pt(1)  # Very tight line spacing to maximize content
            p.line_spacing = 1.5  # Use the same line spacing as text frame
        
        # Apply overflow detection and correction
        print(f"🔍 Checking for overflow on slide with {len(lines)} lines...")
        # Convert available_height (Inches object) to float for calculation
        if hasattr(available_height, 'inches'):
            available_height_value = available_height.inches
        else:
            # Try to extract numeric value from string representation
            height_str = str(available_height)
            if 'inches' in height_str:
                available_height_value = float(height_str.replace('inches', '').strip())
            else:
                available_height_value = 8.3  # Default fallback
        
        print(f"📐 Available height for overflow detection: {available_height_value} inches")
        self.detect_and_fix_overflow(text_frame, available_height_value)
        print("✅ Overflow check complete")
        
        # Add icons and visual elements to enhance the slide
        print("🎨 Adding visual enhancements...")
        self.add_slide_icons_and_visuals(slide, lines, is_title)
        print("✅ Visual enhancements complete")
    
    def process_single_image(self, image_path, operation, params):
        """Process a single image based on the selected operation"""
        try:
            if not os.path.exists(self.output_dir):
                os.makedirs(self.output_dir)
            
            img = cv2.imread(image_path)
            if img is None:
                print(f"Could not read image: {image_path}")
                return False
                
            image_name = os.path.basename(image_path)
            
            if operation == 'splitter':
                # Simple image splitting
                rows_per_page = params.get('rows_per_page', 10)
                detect_questions = params.get('detect_questions', True)
                
                if detect_questions:
                    custom_config = r'-l tur --oem 3 --psm 6'
                    ocr_text = pytesseract.image_to_string(img, config=custom_config)
                    has_questions = contains_question_indicators(ocr_text)
                else:
                    has_questions = False
                if has_questions:
                    # Don't split images with questions
                    cleaned_img = remove_empty_rows_and_columns(img)
                    cleaned_img = add_outer_border(cleaned_img, top_border=50, bottom_border=50)
                    cleaned_img = cv2.resize(cleaned_img, (1920, 1080))
                    output_path = os.path.join(self.output_dir, image_name)
                    cv2.imwrite(output_path, cleaned_img)
                else:
                    # Intelligent split that avoids cutting through text
                    max_splits = params.get('max_splits', 2)  # Default to 2 parts, but configurable
                    
                    if max_splits > 2:
                        # Use advanced multi-split analysis for complex images
                        split_points = self.find_multiple_split_points(img, max_splits)
                    else:
                        # Use simple 2-way split with text analysis
                        split_points = self.find_optimal_split_points(img)
                    
                    print(f"🎯 Applying {len(split_points)} intelligent splits")
                    
                    for i, (start_y, end_y) in enumerate(split_points):
                        crop_img = img[start_y:end_y, :]
                        
                        crop_img = remove_empty_rows_and_columns(crop_img)
                        crop_img = add_outer_border(crop_img, top_border=100, bottom_border=100, left_border=40, right_border=40)
                        if not is_image_blank(crop_img):
                            crop_img = add_outer_border(crop_img, top_border=50, bottom_border=50)
                            crop_img = cv2.resize(crop_img, (1920, 1080))
                            output_path = os.path.join(self.output_dir, f"{image_name}_{i}.png")
                            cv2.imwrite(output_path, crop_img)
            
            elif operation == 'title_cropper':
                # Simple title cropping
                custom_config = r'-l tur --oem 3 --psm 6'
                d = pytesseract.image_to_data(img, config=custom_config, output_type=Output.DICT)
                
                if len(d['text']) > 0:
                    # Find first non-empty text
                    first_text_top = None
                    for i, text in enumerate(d['text']):
                        if text.strip():
                            first_text_top = d['top'][i]
                            break
                    
                    if first_text_top and first_text_top > 50:
                        crop_img = img[first_text_top-20:, :]
                        img = cv2.resize(crop_img, (1920, 1080))
                
                cv2.imwrite(os.path.join(self.output_dir, image_name), img)
            
            elif operation == 'blank_remover':
                # Check if image is blank
                if is_image_blank(img, background_threshold=98.0):
                    blank_dir = os.path.join(self.output_dir, "blank_images")
                    if not os.path.exists(blank_dir):
                        os.makedirs(blank_dir)
                    
                    pil_img = Image.fromarray(cv2.cvtColor(img, cv2.COLOR_BGR2RGB))
                    output_path = os.path.join(blank_dir, image_name)
                    pil_img.save(output_path)
                else:
                    pil_img = Image.fromarray(cv2.cvtColor(img, cv2.COLOR_BGR2RGB))
                    output_path = os.path.join(self.output_dir, image_name)
                    pil_img.save(output_path)
            
            return True
            
        except Exception as e:
            print(f"Error processing image {image_path}: {e}")
            return False
    def process_files_thread(self):
        """Process files in a separate thread"""
        try:
            if not self.input_files:
                messagebox.showwarning(self.lang.get_text('no_input'), 
                    self.lang.get_text('no_input_msg'))
                return
                
            operation = self.operation_var.get()
            
            # Get parameters
            try:
                lines_per_slide = int(self.lines_per_slide_var.get())
                if lines_per_slide < 1:
                    lines_per_slide = 7
            except ValueError:
                lines_per_slide = 7
            
            try:
                rows_per_page = int(self.rows_per_page_var.get())
                if rows_per_page < 1:
                    rows_per_page = 10
            except ValueError:
                rows_per_page = 10
            
            try:
                max_splits = int(self.max_splits_var.get())
                if max_splits < 1:
                    max_splits = 2
                elif max_splits > 5:  # Limit to reasonable number
                    max_splits = 5
            except (ValueError, AttributeError):
                max_splits = 2
            
            params = {
                'lines_per_slide': lines_per_slide,
                'rows_per_page': rows_per_page,
                'detect_questions': self.detect_questions_var.get(),
                'max_splits': max_splits
            }
            
            # Reset progress
            self.progress_var.set(0)
            self.status_var.set(self.lang.get_text('processing'))
            
            total_files = len(self.input_files)
            success_count = 0
            
            if operation == 'doc_to_pptx':
                # Process documents
                for i, doc_path in enumerate(self.input_files):
                    if self.processing_cancelled:
                        break
                        
                    if self.convert_document_to_pptx(doc_path, lines_per_slide):
                        success_count += 1
                    
                    progress = (i + 1) / total_files * 100
                    self.progress_var.set(progress)
                    self.root.update()
                
                self.status_var.set(self.lang.get_text('document_conversion_complete'))
                success_text = f"{self.lang.get_text('successfully_converted')} {success_count} / {total_files}"
                messagebox.showinfo(self.lang.get_text('complete'), f"{success_text} PowerPoint!")
            
            else:
                # Process images
                for i, image_path in enumerate(self.input_files):
                    if self.processing_cancelled:
                        break
                        
                    if self.process_single_image(image_path, operation, params):
                        success_count += 1
                    
                    progress = (i + 1) / total_files * 100
                    self.progress_var.set(progress)
                    status_text = f"{self.lang.get_text('processed_files')} {i + 1} / {total_files}"
                    self.status_var.set(status_text)
                    self.root.update()
                
                if operation == 'blank_remover':
                    blank_dir = os.path.join(self.output_dir, "blank_images")
                    blank_count = len(os.listdir(blank_dir)) if os.path.exists(blank_dir) else 0
                    status_text = f"{self.lang.get_text('processing_complete')} {blank_count} {self.lang.get_text('found_blank_images')}"
                    self.status_var.set(status_text)
                    result_text = f"{blank_count} {self.lang.get_text('blank_images_found')} {total_files}"
                    messagebox.showinfo(self.lang.get_text('complete'), result_text)
                else:
                    self.status_var.set(self.lang.get_text('processing_complete'))
                    success_text = f"{self.lang.get_text('successfully_processed')} {success_count} / {total_files}"
                    messagebox.showinfo(self.lang.get_text('complete'), success_text)
            
        except Exception as e:
            error_text = f"{self.lang.get_text('error_occurred')}{str(e)}"
            messagebox.showerror(self.lang.get_text('error'), error_text)
        finally:
            self.cancel_button.config(state=tk.DISABLED)
            self.processing_cancelled = False
    
    def process_files(self):
        """Start processing files"""
        self.processing_cancelled = False
        self.cancel_button.config(state=tk.NORMAL)
        thread = threading.Thread(target=self.process_files_thread)
        thread.daemon = True
        thread.start()
    def cancel_processing(self):
        """Cancel the current processing"""
        self.processing_cancelled = True
        self.status_var.set(self.lang.get_text('cancelling'))
    def open_output_folder(self):
        """Open the output folder in file explorer"""
        if os.path.exists(self.output_dir):
            os.startfile(self.output_dir)
        else:
            messagebox.showwarning(self.lang.get_text('folder_not_found'), 
                self.lang.get_text('folder_not_found_msg'))
    
    def extract_docx_content_with_tables(self, document_path):
        """Extract content from Word document, preserving tables"""
        content_items = []
        
        try:
            doc = docx.Document(document_path)
            
            # Get all paragraphs and tables in document order
            for element in doc.element.body:
                if element.tag.endswith('p'):  # Paragraph
                    # Create paragraph object from element
                    for para in doc.paragraphs:
                        if para._element is element and para.text.strip():
                            content_items.append({
                                'type': 'text',
                                'content': para.text.strip()
                            })
                            break
                            
                elif element.tag.endswith('tbl'):  # Table
                    # Create table object from element
                    for table in doc.tables:
                        if table._element is element:
                            table_data = []
                            for row in table.rows:
                                row_data = [cell.text.strip() for cell in row.cells]
                                table_data.append(row_data)
                            
                            # Only add non-empty tables
                            if table_data and any(any(cell for cell in row) for row in table_data):
                                content_items.append({
                                    'type': 'table',
                                    'content': table_data,
                                    'rows': len(table_data),
                                    'cols': len(table_data[0]) if table_data else 0
                                })
                            break
                        
        except Exception as e:
            print(f"Error extracting Word content with tables: {e}")
            # Fallback to simple extraction
            try:
                doc = docx.Document(document_path)
                # Get all paragraphs first
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        content_items.append({
                            'type': 'text',
                            'content': paragraph.text.strip()
                        })
                
                # Then get all tables
                for table in doc.tables:
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    
                    if table_data and any(any(cell for cell in row) for row in table_data):                        content_items.append({
                            'type': 'table',
                            'content': table_data,
                            'rows': len(table_data),
                            'cols': len(table_data[0]) if table_data else 0
                        })
            except Exception as fallback_error:
                print(f"Fallback extraction also failed: {fallback_error}")
        
        return content_items
    
    def extract_pdf_content_with_tables(self, document_path):
        """Extract content from PDF with high-quality text extraction for Turkish characters"""
        content_items = []
        
        # Method 1: Try PyMuPDF first (best for Turkish text)
        if PYMUPDF_AVAILABLE:
            print("🚀 Using PyMuPDF for high-quality Turkish text extraction...")
            pymupdf_result = self.extract_text_with_pymupdf(document_path)
            if pymupdf_result:
                return pymupdf_result
        
        # Method 2: Fall back to pdfplumber for table detection
        try:
            try:
                import pdfplumber
                use_pdfplumber = True
                print("📖 Using pdfplumber for enhanced table detection")
            except ImportError:
                use_pdfplumber = False
                print("pdfplumber not available, using basic PDF extraction")
            
            if use_pdfplumber:
                with pdfplumber.open(document_path) as pdf:
                    for page_num, page in enumerate(pdf.pages):
                        # Extract tables first
                        tables = page.extract_tables()
                        if tables:
                            for table in tables:
                                # Filter empty rows and cells
                                cleaned_table = []
                                for row in table:
                                    if row and any(cell and str(cell).strip() for cell in row):
                                        cleaned_row = [str(cell).strip() if cell else '' for cell in row]
                                        cleaned_table.append(cleaned_row)
                                
                                if cleaned_table:
                                    content_items.append({
                                        'type': 'table',
                                        'content': cleaned_table,
                                        'rows': len(cleaned_table),
                                        'cols': len(cleaned_table[0]) if cleaned_table else 0,
                                        'page': page_num + 1
                                    })
                        
                        # Extract text
                        text = page.extract_text()
                        if text:
                            lines = [line.strip() for line in text.split('\n') if line.strip()]
                            for line in lines:
                                # Simple heuristic to avoid duplicating table content
                                if not self.is_likely_table_text(line):
                                    content_items.append({
                                        'type': 'text',
                                        'content': self.clean_extracted_text(line),
                                        'page': page_num + 1
                                    })
            else:
                # Fallback to PyPDF2 with basic table detection
                with open(document_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page_num, page in enumerate(pdf_reader.pages):
                        text = page.extract_text()
                        if text:
                            lines = [line.strip() for line in text.split('\n') if line.strip()]
                            for line in lines:
                                # Simple table detection based on patterns
                                if self.is_likely_table_text(line):
                                    # Try to parse as table row
                                    table_row = self.parse_table_row(line)
                                    if table_row:
                                        content_items.append({
                                            'type': 'table',
                                            'content': [table_row],
                                            'rows': 1,
                                            'cols': len(table_row),
                                            'page': page_num + 1
                                        })
                                    else:
                                        content_items.append({
                                            'type': 'text',
                                            'content': self.clean_extracted_text(line),
                                            'page': page_num + 1
                                        })
                                else:
                                    content_items.append({
                                        'type': 'text',
                                        'content': self.clean_extracted_text(line),
                                        'page': page_num + 1
                                    })
                                    
        except Exception as e:
            print(f"Error extracting PDF content: {e}")            # Ultimate fallback
            try:
                with open(document_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page_num, page in enumerate(pdf_reader.pages):
                        text = page.extract_text()
                        if text:
                            lines = [line.strip() for line in text.split('\n') if line.strip()]
                            for line in lines:                                content_items.append({
                                    'type': 'text',
                                    'content': self.clean_extracted_text(line),
                                    'page': page_num + 1
                                })
            except Exception as final_error:
                print(f"All PDF extraction methods failed: {final_error}")
        
        return content_items
    
    def reconstruct_line_from_spans(self, sorted_spans):
        """Reconstruct a text line from sorted spans, handling spacing intelligently"""
        if not sorted_spans:
            return ""
        
        line_text = ""
        prev_x_end = None
        
        for span in sorted_spans:
            span_text = span['text']
            bbox = span['bbox']
            x_start = bbox[0]
            x_end = bbox[2]
            
            # Add appropriate spacing between spans
            if prev_x_end is not None:
                gap = x_start - prev_x_end
                
                # Determine if we need to add space based on gap size
                if gap > 5:  # Significant gap - likely a space or tab
                    if not line_text.endswith(' ') and not span_text.startswith(' '):
                        line_text += ' '
                elif gap < -2:  # Overlapping - might be bold/italic overlay
                    # Don't add space for overlapping spans
                    pass
                # For small gaps (0-5), rely on existing spacing in the text
            
            line_text += span_text
            prev_x_end = x_end
        
        return line_text
    
    def group_spans_by_visual_lines(self, spans_data, y_tolerance=2):
        """Group spans by their visual line position with tolerance for slight y variations"""
        line_groups = {}
        
        for span_info in spans_data:
            y_pos = span_info['bbox'][1]
            
            # Find existing line group within tolerance
            matched_y = None
            for existing_y in line_groups.keys():
                if abs(y_pos - existing_y) <= y_tolerance:
                    matched_y = existing_y
                    break
            
            # Use existing group or create new one
            target_y = matched_y if matched_y is not None else y_pos
            
            if target_y not in line_groups:
                line_groups[target_y] = []
            
            line_groups[target_y].append(span_info)
        
        return line_groups
    
    def extract_text_with_pymupdf(self, document_path):
        """Extract text using PyMuPDF (fitz) for high-quality Turkish text extraction"""
        content_items = []
        
        try:
            if not PYMUPDF_AVAILABLE:
                return None
                
            doc = fitz.open(document_path)
            
            for page_num in range(doc.page_count):
                page = doc[page_num]
                
                # Method 1: Try structured text extraction first (preserves formatting)
                text_dict = page.get_text("dict")
                  # Method 2: Also get raw text for comparison and fallback
                raw_text = page.get_text()
                
                current_paragraph = []
                extracted_blocks = []
                
                # Collect all spans with their positions
                all_spans = []
                
                # Process structured text first
                for block in text_dict["blocks"]:
                    if "lines" in block:  # Text block
                        for line in block["lines"]:
                            for span in line["spans"]:
                                span_text = span["text"]
                                if span_text.strip():
                                    all_spans.append({
                                        'text': span_text,
                                        'bbox': span["bbox"],
                                        'x': span["bbox"][0],
                                        'y': span["bbox"][1]
                                    })
                
                # Group spans by visual lines (y-coordinate with tolerance)
                line_groups = self.group_spans_by_visual_lines(all_spans, y_tolerance=3)
                
                # Sort lines by y-coordinate and process each visual line
                sorted_y_positions = sorted(line_groups.keys())
                
                for y_pos in sorted_y_positions:
                    spans_in_line = line_groups[y_pos]
                    
                    # Sort spans within the line by x-coordinate (left to right)
                    sorted_spans = sorted(spans_in_line, key=lambda s: s['x'])
                      # Combine spans into a single line, handling spacing intelligently
                    line_text = self.reconstruct_line_from_spans(sorted_spans)
                    
                    # Clean and add the reconstructed line
                    if line_text.strip():
                        extracted_blocks.append(line_text.strip())
                
                print(f"Extracted {len(extracted_blocks)} visual lines from page {page_num + 1}")
                if extracted_blocks:
                    print("Sample lines:", extracted_blocks[:3])
                
                # If structured extraction worked, use it
                if extracted_blocks:
                    for block_text in extracted_blocks:
                        # Split into logical units (sentences or paragraphs)
                        sentences = self.split_into_logical_units(block_text)
                        for sentence in sentences:
                            sentence = sentence.strip()
                            if sentence:
                                # Check if this looks like a table row
                                if self.is_likely_table_text(sentence):
                                    table_row = self.parse_table_row(sentence)
                                    if table_row:
                                        content_items.append({
                                            'type': 'table',
                                            'content': [table_row],
                                            'rows': 1,
                                            'cols': len(table_row),
                                            'page': page_num + 1
                                        })
                                    else:
                                        content_items.append({
                                            'type': 'text',
                                            'content': self.minimal_text_cleanup(sentence),
                                            'page': page_num + 1
                                        })
                                else:
                                    content_items.append({
                                        'type': 'text',
                                        'content': self.minimal_text_cleanup(sentence),
                                        'page': page_num + 1
                                    })
                else:
                    # Fallback to raw text if structured extraction failed
                    if raw_text.strip():
                        lines = [line.strip() for line in raw_text.split('\n') if line.strip()]
                        for line in lines:
                            if self.is_likely_table_text(line):
                                table_row = self.parse_table_row(line)
                                if table_row:
                                    content_items.append({
                                        'type': 'table',
                                        'content': [table_row],
                                        'rows': 1,
                                        'cols': len(table_row),
                                        'page': page_num + 1
                                    })
                                else:
                                    content_items.append({
                                        'type': 'text',
                                        'content': self.minimal_text_cleanup(line),
                                        'page': page_num + 1
                                    })
                            else:
                                content_items.append({
                                    'type': 'text',
                                    'content': self.minimal_text_cleanup(line),
                                    'page': page_num + 1
                                })
            
            doc.close()
            print(f"✅ PyMuPDF extracted {len(content_items)} items with high precision")
            return content_items
            
        except Exception as e:
            print(f"PyMuPDF extraction failed: {e}")
            return None
    
    def split_into_logical_units(self, text):
        """Split text into logical units (sentences, bullet points, etc.)"""
        # Split on sentence endings, but be careful with abbreviations
        import re
        
        # Common abbreviations that shouldn't end sentences
        abbreviations = ['Dr', 'Mr', 'Mrs', 'Ms', 'Prof', 'vs', 'etc', 'Inc', 'Ltd', 'Co']
        
        # Split on periods, exclamation marks, question marks followed by space and capital letter
        sentences = re.split(r'(?<=[.!?])\s+(?=[A-ZÜĞŞÇÖI])', text)
        
        # Also split on bullet points and numbered lists
        result = []
        for sentence in sentences:
            # Check for bullet points or numbered lists
            bullet_parts = re.split(r'(?:^|\s)([•·▪▫‣⁃]\s*|\d+\.\s*|[a-zA-Z]\.\s*)', sentence)
            if len(bullet_parts) > 1:
                for part in bullet_parts:
                    if part and part.strip():
                        result.append(part.strip())
            else:
                result.append(sentence.strip())
        return [s for s in result if s]
    
    def minimal_text_cleanup(self, text):
        """Minimal text cleanup for PyMuPDF extracted text (which should be cleaner)"""
        if not text:
            return text
        
        # Basic whitespace normalization
        text = ' '.join(text.split())
        
        # Only fix the most common spacing issues around punctuation
        import re
        text = re.sub(r'\s+([,.!?;:])', r'\1', text)
        text = re.sub(r'(\()\s+', r'\1', text)
        text = re.sub(r'\s+(\))', r'\1', text)
        text = re.sub(r'(\")\s+', r'\1', text)
        text = re.sub(r'\s+(\")', r'\1', text)
        
        # Apply targeted Turkish character fixes for common patterns
        turkish_fixes = {
            # Most critical Turkish character spacing issues
            'ş tirme': 'ştirme', 'ş tir': 'ştir', 'ş me': 'şme', 'ş ma': 'şma',
            'ğ ın': 'ğın', 'ğ ına': 'ğına', 'ğ e': 'ğe', 'ğ a': 'ğa', 'ğ u': 'ğu',
            'ç ın': 'çın', 'ç ına': 'çına', 'ç e': 'çe', 'ç a': 'ça', 'ç i': 'çi',
            'ü ş': 'üş', 'ü r': 'ür', 'ü n': 'ün', 'ü m': 'üm',
            'ö n': 'ön', 'ö r': 'ör', 'ö l': 'öl', 'ö z': 'öz',
            
            # Common pattern fixes for both Turkish and English
            'i ş': 'iş', 'i n': 'in', 'i r': 'ir', 'i l': 'il',
            't he': 'the', 'a nd': 'and', 'f or': 'for', 't o': 'to'
        }
        
        for broken, fixed in turkish_fixes.items():
            if broken in text:
                text = text.replace(broken, fixed)
        
        return text
    
    def is_likely_table_text(self, text):
        """Heuristic to detect if text line is likely from a table"""
        # Look for common table patterns
        patterns = [
            r'\s+\|\s+',              # Pipe separators
            r'\t{2,}',                # Multiple tabs
            r'\s{4,}',                # Multiple spaces (4 or more)
            r'^\s*\d+\.\s+.*\s+\d+',  # Number, text, number pattern
            r'.*\s+\$\d+',            # Text followed by dollar amount
            r'.*\s+\d+%',             # Text followed by percentage
        ]
        
        for pattern in patterns:
            if re.search(pattern, text):
                return True
                
        # Check if line has structure like "Item1    Item2    Item3"
        parts = text.split()
        if len(parts) >= 3:
            # Check for significant spacing between words
            original_spaces = len(text) - len(text.replace(' ', ''))
            if original_spaces > len(parts) * 3:  # More than 3 spaces per word on average
                return True
                
        return False
    
    def parse_table_row(self, text):
        """Attempt to parse a text line as a table row"""
        # Try different separators
        separators = ['|', '\t', '    ', '   ']  # Tab, multiple spaces
        
        for sep in separators:
            if sep in text:
                cells = [cell.strip() for cell in text.split(sep)]
                # Filter out empty cells and ensure we have at least 2 columns
                cells = [cell for cell in cells if cell]
                if len(cells) >= 2:
                    return cells
        
        return None

    def create_table_slide(self, prs, slide_layout, table_data):
        """Create a slide with a table"""
        slide = prs.slides.add_slide(slide_layout)
        
        # Calculate table dimensions
        rows = len(table_data)
        cols = max(len(row) for row in table_data) if table_data else 1
        
        # Position table in center of slide with appropriate sizing
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(15)    # Almost full width
        height = Inches(6.5)  # Leave space for potential title
        
        # Add table shape
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        # Populate table
        for row_idx, row_data in enumerate(table_data):
            for col_idx, cell_data in enumerate(row_data):
                if col_idx < cols:  # Ensure we don't exceed table columns
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(cell_data) if cell_data else ''
                    
                    # Style the cells
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.name = 'Calibri'
                        paragraph.font.size = Pt(11)  # Slightly smaller for tables
                          # Make header row bold and centered
                        if row_idx ==  0:
                            paragraph.font.bold = True
                            paragraph.alignment = PP_ALIGN.CENTER
                            cell.fill.solid()
                            accent_rgb = self.get_color_setting('accent_color', '#4472c4')
                            cell.fill.fore_color.rgb = RGBColor(*accent_rgb)
                            paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
                        else:
                            paragraph.font.bold = False
                            paragraph.alignment = PP_ALIGN.LEFT
        
        # Apply table styling
        try:
            table.style = 'LightShading-Accent1'  # Built-in table style
        except:
            pass  # Style might not be available
        
        return slide

    def add_slide_icons_and_visuals(self, slide, lines, is_title):
        """Add icons and visual elements to enhance slide appearance"""
        try:
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.dml.color import RGBColor
            
            # Add a decorative header line for titles
            if is_title and len(lines) > 0:
                self.add_decorative_header_line(slide)
            
            # Add content-based icons
            self.add_content_icons(slide, lines)
            
        except Exception as e:
            print(f"⚠️  Note: Could not add visual elements: {e}")
    
    def add_decorative_header_line(self, slide):
        """Add a decorative line under the title"""
        try:
            from pptx.enum.shapes import MSO_SHAPE
            
            # Position the line below the title area
            # Assuming title takes about 0.8-1.0 inches, place line below it
            left = Inches(2)
            top = Inches(1.8)  # Moved down to be below title text
            width = Inches(12)
            height = Inches(0.02)
            line_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, width, height
            )
            line_shape.fill.solid()
            accent_rgb = self.get_color_setting('accent_color', '#1f497d')
            line_shape.fill.fore_color.rgb = RGBColor(*accent_rgb)
            line_shape.line.color.rgb = RGBColor(*accent_rgb)
            
        except Exception as e:
            print(f"⚠️  Could not add header line: {e}")
    
    def add_content_icons(self, slide, lines):
        """Add icons based on content type"""
        try:
            from pptx.enum.shapes import MSO_SHAPE
            
            icon_x = Inches(0.1)  # Left margin for icons
            start_y = Inches(1.5)  # Start below title area
            
            for i, line in enumerate(lines):
                if i == 0:  # Skip title line
                    continue
                    
                y_position = start_y + Inches(i * 0.3)  # Space icons vertically
                
                # Determine icon type based on content
                icon_shape = None
                icon_color = RGBColor(68, 84, 106)  # Default blue-gray
                
                if any(keyword in line.lower() for keyword in ['question', '?', 'soru']):
                    # Question mark icon for questions
                    icon_shape = MSO_SHAPE.OVAL
                    icon_color = RGBColor(255, 193, 7)  # Amber
                # elif line.startswith(('•', '-', '*', '▪', '○')) or any(line.startswith(f'{j}.') for j in range(1, 10)):
                #     # Bullet point icon
               
                #     icon_shape = MSO_SHAPE.OVAL
                #     icon_color = RGBColor(40, 167, 69)  # Green
                elif any(keyword in line.lower() for keyword in ['important', 'note', 'warning', 'dikkat']):
                    # Important/warning icon
                    icon_shape = MSO_SHAPE.PENTAGON
                    icon_color = RGBColor(220, 53, 69)  # Red
                elif any(keyword in line.lower() for keyword in ['example', 'örnek', 'sample']):
                    # Example icon
                    icon_shape = MSO_SHAPE.HEXAGON
                    icon_color = RGBColor(111, 66, 193)  # Purple
                
                # Add the icon if we determined one
                if icon_shape and i < 15:  # Limit icons to prevent clutter
                    self.add_small_icon(slide, icon_x, y_position, icon_shape, icon_color)
                    
        except Exception as e:
            print(f"⚠️  Could not add content icons: {e}")
    
    def add_small_icon(self, slide, x, y, shape_type, color):
        """Add a small icon at the specified position"""
        try:
            size = Inches(0.15)  # Small icon size
            
            icon = slide.shapes.add_shape(shape_type, x, y, size, size)
            icon.fill.solid()
            icon.fill.fore_color.rgb = color
            icon.line.color.rgb = color
            
        except Exception as e:
            print(f"⚠️  Could not add icon: {e}")
    
    def is_better_title(self, line, line_index=0, context_lines=None):
        """
        Improved title detection logic with stricter criteria:
        1. Must be short (< 80 chars)
        2. Should be uppercase OR contain title keywords OR have strong title structure
        3. Should NOT end with punctuation that suggests continuation
        4. Should NOT be obviously a question in the middle of content
        5. Should NOT contain common non-title patterns (verb endings, articles, etc.)
        """
        if not line or len(line.strip()) == 0:
            return False
            
        line = line.strip()
        
        # Too long to be a title
        if len(line) > 80:
            return False
            
        # If it ends with continuation punctuation, probably not a title
        if line.endswith(('...', ':', ';', ',', '-')):
            return False
            
        # If it's a question in the middle of content, probably not a title
        if line.endswith('?') and line_index > 0:
            # Check if this looks like a content question vs title question
            question_words = ['ne', 'nasıl', 'neden', 'niçin', 'kim', 'nerede', 'ne zaman',
                             'where', 'what', 'how', 'why', 'when', 'who', 'which']
            if any(word in line.lower() for word in question_words):
                return False  # Content question, not title
                  # Negative indicators - patterns that suggest this is NOT a title
        line_lower = line.lower()
        line_words = line_lower.split()
        
        # Common Turkish verb endings and patterns that indicate regular content
        # Check for whole words and specific endings
        non_title_words = ['bir', 'bu', 'şu', 'o', 'ile', 'için', 'gibi', 
                          'the', 'a', 'an', 'is', 'are', 'was', 'were']
        non_title_endings = ['dir', 'dır', 'dur', 'dür', 'ler', 'lar', 
                            'den', 'dan', 'ten', 'tan', 'nın', 'nin', 'nun', 'nün',
                            'da', 'de', 'ta', 'te']
        specific_verbs = ['yetkilidir', 'edilir', 'yapılır', 'olur', 'bulur']
          # Check for negative patterns
        has_non_title_pattern = False
        
        # Check whole words first (most definitive)
        if any(word in line_words for word in non_title_words):
            has_non_title_pattern = True
        
        # Check specific verb forms (but allow title case to override)
        elif any(verb in line_lower for verb in specific_verbs):
            # Check if it's title case first
            if len(line_words) <= 3:
                capitalized_words = [w for w in line.split() if w and w[0].isupper()]
                if len(capitalized_words) >= len(line.split()) * 0.6:  # 60% capitalized
                    has_non_title_pattern = False  # Title case overrides verb endings
                else:
                    has_non_title_pattern = True
            else:
                has_non_title_pattern = True
                
        # Check word endings (but be more careful)
        elif any(word.endswith(ending) for word in line_words for ending in non_title_endings):
            # Additional check: if it's a single word ending with these, it might still be a title
            # if it's capitalized and not obviously a verb
            if len(line_words) == 1 and line_words[0][0].isupper():
                has_non_title_pattern = False  # Single capitalized word might be a title
            # Also check if it's clearly title case (most words capitalized)
            elif len(line_words) <= 3:
                capitalized_words = [w for w in line.split() if w and w[0].isupper()]
                if len(capitalized_words) >= len(line.split()) * 0.6:  # 60% capitalized
                    has_non_title_pattern = False  # Title case overrides endings
                else:
                    has_non_title_pattern = True
            else:
                has_non_title_pattern = True
        
        # Strong title indicators
        if line.isupper() and len(line) > 2:
            return True
            
        # Title keywords (case insensitive)
        title_keywords = ['başlık', 'giriş', 'bölüm', 'sonuç', 'özet', 'genel bakış', 
                         'introduction', 'chapter', 'conclusion', 'summary', 'overview',
                         'madde', 'kısım', 'section', 'part']
        has_title_keyword = any(keyword in line_lower for keyword in title_keywords)
        
        if has_title_keyword:
            return True
            
        # Numbers at start might indicate sections/chapters
        if line_lower.startswith(('1.', '2.', '3.', '4.', '5.', '6.', '7.', '8.', '9.', 
                                 'i.', 'ii.', 'iii.', 'iv.', 'v.', 'a.', 'b.', 'c.')):
            return True
              # For short text at the beginning, be more restrictive
        if line_index == 0 and len(line) < 50:
            # Only consider it a title if it doesn't have non-title patterns
            # AND has some title-like characteristics
            if has_non_title_pattern:
                return False
                
            # Additional title indicators for short text
            # Check if it looks like a proper noun or title case
            words = line.split()
            if len(words) <= 4:  # Short phrases might be titles
                # Check if most words are capitalized (title case)
                capitalized_words = [w for w in words if w and w[0].isupper()]
                if len(capitalized_words) >= len(words) * 0.6:  # At least 60% capitalized
                    return True
                    
            # Very short (< 20 chars), no negative patterns, and meaningful length might be title
            if len(line) < 20 and not has_non_title_pattern and len(line) > 3:
                return True
                
        return False
    
    def add_unicode_symbols_to_text(self, text, is_title_line=False):
        """Enhanced Unicode symbols with better logic for when to add them"""
        # Don't add question mark symbol to titles 
        if is_title_line:
            return text
            
        # Only add question symbol to actual questions (not just any text with ?)
        if text.strip().endswith('?'):            # Check if it's actually a question vs just ends with ?
            question_indicators = ['ne', 'nasıl', 'neden', 'niçin', 'kim', 'nerede', 'ne zaman',
                                 'what', 'how', 'why', 'when', 'where', 'who', 'which']
            if any(indicator in text.lower() for indicator in question_indicators):
                return f"❓ {text}"
        elif any(keyword in text.lower() for keyword in ['important', 'önemli', 'dikkat']):
            return f"⚠️ {text}"
        elif any(keyword in text.lower() for keyword in ['example', 'örnek']):
            return f"💡 {text}"
        elif any(keyword in text.lower() for keyword in ['note', 'not']):
            return f"📝 {text}"
        elif text.startswith(('•', '-', '*')):
            return f"▶️ {text[1:].strip()}"  # Replace bullet with arrow
        else:
            return text
    
    def on_closing(self):
        """Handle application closing"""
        try:
            self.save_current_values()
        except:
            pass
        self.root.destroy()
    
    def load_saved_values(self):
        """Load saved configuration values and restore UI state"""
        try:
            if os.path.exists("app_settings.json"):
                with open("app_settings.json", "r", encoding="utf-8") as f:
                    settings = json.load(f)
                      # Restore UI state
                    if hasattr(self, 'operation_var') and 'operation' in settings:
                        self.operation_var.set(settings['operation'])
                    
                    if hasattr(self, 'lines_per_slide_var') and 'lines_per_slide' in settings:
                        self.lines_per_slide_var.set(settings['lines_per_slide'])
                    
                    if hasattr(self, 'line_spacing_var') and 'line_spacing' in settings:
                        self.line_spacing_var.set(settings['line_spacing'])
                    
                    if hasattr(self, 'rows_per_page_var') and 'rows_per_page' in settings:
                        self.rows_per_page_var.set(settings['rows_per_page'])
                    
                    if hasattr(self, 'detect_questions_var') and 'detect_questions' in settings:
                        self.detect_questions_var.set(settings['detect_questions'])
                    
                    # Restore output directory
                    if 'output_directory' in settings:
                        self.output_dir = settings['output_directory']
                        if hasattr(self, 'output_dir_label'):
                            self.output_dir_label.config(text=self.output_dir)
                    
        except Exception as e:
            print(f"Could not load saved values: {e}")
    
    def save_current_values(self):
        """Save current configuration values including UI state"""
        try:
            settings = {
                "last_saved": time.time(),
                # UI state values
                "language": self.lang.get_current_language(),
                "operation": self.operation_var.get() if hasattr(self, 'operation_var') else "splitter",                "lines_per_slide": self.lines_per_slide_var.get() if hasattr(self, 'lines_per_slide_var') else "7",
                "line_spacing": self.line_spacing_var.get() if hasattr(self, 'line_spacing_var') else "1.5",
                "rows_per_page": self.rows_per_page_var.get() if hasattr(self, 'rows_per_page_var') else "10",
                "detect_questions": self.detect_questions_var.get() if hasattr(self, 'detect_questions_var') else True,
                # Output directory
                "output_directory": self.output_dir if hasattr(self, 'output_dir') else os.path.join(os.getcwd(), "cropped"),                # Window settings
                "window_geometry": f"{self.root.winfo_width()}x{self.root.winfo_height()}" if hasattr(self, 'root') else "800x900",
                "window_position": f"+{self.root.winfo_x()}+{self.root.winfo_y()}" if hasattr(self, 'root') and self.root.winfo_x() >= 0 and self.root.winfo_y() >= 0 else None
            }
            with open("app_settings.json", "w", encoding="utf-8") as f:
                json.dump(settings, f, indent=2, ensure_ascii=False)
        except Exception as e:
            print(f"Could not save values: {e}")
    
    def on_language_change(self, event=None):

        """Handle language change"""
        try:
            selected_language = self.language_var.get()
            if selected_language != self.lang.get_current_language():
                self.lang.set_language(selected_language)
                self.update_all_ui_text()
        except Exception as e:
            print(f"Language change error: {e}")
    
    def update_language_combo(self, combo):
        """Update language combo box"""
        try:
            # Set the display text based on the current language
            current_lang = self.lang.get_current_language()
            combo.set(current_lang)
            
            # Update the values to show proper display names
            display_values = []
            for lang_code in ['tr', 'en']:
                if lang_code == 'tr':
                    display_values.append('tr')
                else:
                    display_values.append('en')
            combo['values'] = display_values
        except Exception as e:
            print(f"Could not update language combo: {e}")
    
    def update_all_ui_text(self):
        """Update all UI elements with new language text"""
        try:
            # Update window title
            self.root.title(self.lang.get_text('window_title'))
              # Update all widgets that have text
            if hasattr(self, 'widgets'):
                # Update frame labels
                if 'input_frame' in self.widgets:
                    self.widgets['input_frame'].config(text=self.lang.get_text('input_section'))
                if 'output_frame' in self.widgets:
                    self.widgets['output_frame'].config(text=self.lang.get_text('output_section'))
                if 'operation_frame' in self.widgets:
                    self.widgets['operation_frame'].config(text=self.lang.get_text('operation_section'))
                if 'param_frame' in self.widgets:
                    self.widgets['param_frame'].config(text=self.lang.get_text('parameters_section'))
                if 'preview_frame' in self.widgets:
                    self.widgets['preview_frame'].config(text=self.lang.get_text('preview_section'))
            
                # Update buttons from widgets dictionary
                if 'select_folder_button' in self.widgets:
                    self.widgets['select_folder_button'].config(text=self.lang.get_text('select_folder'))
                if 'select_images_button' in self.widgets:
                    self.widgets['select_images_button'].config(text=self.lang.get_text('select_images'))
                if 'select_docs_button' in self.widgets:
                    self.widgets['select_docs_button'].config(text=self.lang.get_text('select_documents'))
                if 'output_button' in self.widgets:
                    self.widgets['output_button'].config(text=self.lang.get_text('select_output_folder'))
                if 'process_button' in self.widgets:
                    self.widgets['process_button'].config(text=self.lang.get_text('process'))
                if 'cancel_button' in self.widgets:
                    self.widgets['cancel_button'].config(text=self.lang.get_text('cancel'))
                if 'open_folder_button' in self.widgets:
                    self.widgets['open_folder_button'].config(text=self.lang.get_text('open_output_folder'))
                
                # Update parameter labels
                if 'lines_label' in self.widgets:
                    self.widgets['lines_label'].config(text=self.lang.get_text('lines_per_slide'))
                if 'rec_label' in self.widgets:
                    self.widgets['rec_label'].config(text=self.lang.get_text('lines_recommendation'))
                if 'rows_label' in self.widgets:
                    self.widgets['rows_label'].config(text=self.lang.get_text('rows_per_page'))
                if 'question_checkbox' in self.widgets:
                    self.widgets['question_checkbox'].config(text=self.lang.get_text('detect_questions'))
            
            # Update input files label
            if hasattr(self, 'input_files_label'):
                current_text = self.input_files_label.cget('text')
                # Only update if it's showing the default "no files selected" message
                if 'No files' in current_text or 'Dosya' in current_text:
                    self.input_files_label.config(text=self.lang.get_text('no_files_selected'))
            
            # Update operation radio buttons
            if hasattr(self, 'operation_radios'):
                for radio, text_key in self.operation_radios:
                    radio.config(text=self.lang.get_text(text_key))
            
            # Update status
            if hasattr(self, 'status_var'):
                current_status = self.status_var.get()
                # Only update if it's the default ready status
                if 'Ready' in current_status or 'Hazır' in current_status:
                    self.status_var.set(self.lang.get_text('ready'))
            
            # Update dynamic field text
            self.update_dynamic_field_text()
            
            print(f"Language changed to: {self.lang.get_current_language()}")
            
        except Exception as e:
            print(f"Error updating UI text: {e}")
            import traceback
            traceback.print_exc()
    
    def clean_extracted_text(self, text):
        """Clean text extracted from PDFs to fix common spacing and encoding issues"""
        if not text:
            return text
            
        # Remove excessive whitespace and normalize spaces
        text = ' '.join(text.split())
          # Fix common Turkish character spacing issues
        turkish_fixes = {
            # Common Turkish characters that get separated
            's ı': 'sı', 'ş ı': 'şı', 'ş t': 'şt', 'ş m': 'şm', 'ş r': 'şr',
            'ğ ı': 'ğı', 'ğ e': 'ğe', 'ğ a': 'ğa', 'ğ u': 'ğu',
            'ç ı': 'çı', 'ç e': 'çe', 'ç a': 'ça', 'ç i': 'çi',
            'ü ş': 'üş', 'ü r': 'ür', 'ü n': 'ün', 'ü m': 'üm',
            'ö n': 'ön', 'ö r': 'ör', 'ö l': 'öl', 'ö z': 'öz',
            
            # Common word patterns that get broken
            'y e': 'ye', 'y a': 'ya', 'y i': 'yi', 'y u': 'yu', 'y ö': 'yö',
            'd e': 'de', 'd a': 'da', 'd i': 'di', 'd u': 'du',
            't e': 'te', 't a': 'ta', 't i': 'ti', 't u': 'tu',
            'k e': 'ke', 'k a': 'ka', 'k i': 'ki', 'k u': 'ku',
            'l e': 'le', 'l a': 'la', 'l i': 'li', 'l u': 'lu',
            'm e': 'me', 'm a': 'ma', 'm i': 'mi', 'm u': 'mu',
            'n e': 'ne', 'n a': 'na', 'n i': 'ni', 'n u': 'nu',
            'r e': 're', 'r a': 'ra', 'r i': 'ri', 'r u': 'ru',
            'g e': 'ge', 'g a': 'ga', 'g i': 'gi', 'g u': 'gu', 'g ö': 'gö',
            'b e': 'be', 'b a': 'ba', 'b i': 'bi', 'b u': 'bu',
            'v e': 've', 'v a': 'va', 'v i': 'vi', 'v u': 'vu',
            
            # Specific problematic patterns from the example
            'ş tirme': 'ştirme', 'ş tir': 'ştir',
            'ş me': 'şme', 'ş ma': 'şma',
            'ğ ın': 'ğın', 'ğ ına': 'ğına',
            'ç ın': 'çın', 'ç ına': 'çına',
            
            # Common English patterns that might get broken too
            'i on': 'ion', 'i ng': 'ing', 't he': 'the', 'a nd': 'and',
            'o r': 'or', 'i n': 'in', 'f or': 'for', 't o': 'to'
        }
        
        # Apply Turkish fixes
        for broken, fixed in turkish_fixes.items():
            text = text.replace(broken, fixed)
        
        # Fix multiple spaces that might remain
        text = ' '.join(text.split())
        
        # Remove spaces before punctuation
        import re
        text = re.sub(r'\s+([,.!?;:])', r'\1', text)
        
        # Fix spaces after opening and before closing quotes/parentheses
        text = re.sub(r'(\()\s+', r'\1', text)
        text = re.sub(r'\s+(\))', r'\1', text)
        text = re.sub(r'(\")\s+', r'\1', text)
        text = re.sub(r'\s+(\")', r'\1', text)
        
        return text

    def create_dynamic_fields(self, parent_frame):
        """Create all dynamic fields that can be shown/hidden based on operation"""
        
        # PowerPoint conversion fields
        self.create_ppt_conversion_fields(parent_frame)
        
        # Splitter fields  
        self.create_splitter_fields(parent_frame)
        
        # Common fields used by multiple operations
        self.create_common_fields(parent_frame)
    
    def create_ppt_conversion_fields(self, parent_frame):
        """Create fields specific to PowerPoint conversion"""
        
        # Lines per slide
        lines_frame = ttk.Frame(parent_frame)
        lines_label = ttk.Label(lines_frame, text=self.lang.get_text('lines_per_slide'))
        lines_label.pack(anchor=tk.W)
        
        self.lines_per_slide_var = tk.StringVar(value="7")
        lines_entry = ttk.Entry(lines_frame, textvariable=self.lines_per_slide_var, width=10)
        lines_entry.pack(anchor=tk.W, pady=5)
        
        lines_rec_label = ttk.Label(lines_frame, text=self.lang.get_text('lines_recommendation'), font=("Arial", 8))
        lines_rec_label.pack(anchor=tk.W)
        
        self.dynamic_widgets['ppt_lines'] = {
            'frame': lines_frame,
            'widgets': [lines_label, lines_entry, lines_rec_label]
        }
        
        # Line spacing
        spacing_frame = ttk.Frame(parent_frame)
        spacing_label = ttk.Label(spacing_frame, text=self.lang.get_text('line_spacing'))
        spacing_label.pack(anchor=tk.W, pady=(10, 0))
        
        self.line_spacing_var = tk.StringVar(value="1.5")
        spacing_entry = ttk.Entry(spacing_frame, textvariable=self.line_spacing_var, width=10)
        spacing_entry.pack(anchor=tk.W, pady=5)
        
        spacing_rec_label = ttk.Label(spacing_frame, text=self.lang.get_text('line_spacing_recommendation'), font=("Arial", 8))
        spacing_rec_label.pack(anchor=tk.W)
        self.dynamic_widgets['ppt_spacing'] = {
            'frame': spacing_frame,
            'widgets': [spacing_label, spacing_entry, spacing_rec_label]
        }
        
        # Color settings
        self.create_color_selection_fields(parent_frame)
    
    def create_color_selection_fields(self, parent_frame):
        """Create color selection fields for PowerPoint customization"""
        # Title text color
        title_color_frame = ttk.Frame(parent_frame)
        title_color_label = ttk.Label(title_color_frame, text=self.lang.get_text('title_text_color'))
        title_color_label.pack(anchor=tk.W, pady=(10, 0))
        title_color_button = ttk.Button(
            title_color_frame, 
            text="Color: " + (self.settings.get('title_text_color') or '#000000'),
            width=15,
            command=lambda: self.choose_color('title_text_color', title_color_button)
        )
        title_color_button.pack(anchor=tk.W, pady=(5, 0))
        
        self.dynamic_widgets['ppt_title_color'] = {
            'frame': title_color_frame,
            'widgets': [title_color_label, title_color_button],
            'color_button': title_color_button
        }
        
        # Background color
        bg_color_frame = ttk.Frame(parent_frame)
        bg_color_label = ttk.Label(bg_color_frame, text=self.lang.get_text('background_color'))
        bg_color_label.pack(anchor=tk.W, pady=(10, 0))
        bg_color_button = ttk.Button(
            bg_color_frame, 
            text="Color: " + (self.settings.get('background_color') or '#ffffff'),
            width=15,
            command=lambda: self.choose_color('background_color', bg_color_button)
        )
        bg_color_button.pack(anchor=tk.W, pady=(5, 0))
        
        self.dynamic_widgets['ppt_bg_color'] = {
            'frame': bg_color_frame,
            'widgets': [bg_color_label, bg_color_button],
            'color_button': bg_color_button
        }
        
        # Accent color
        accent_color_frame = ttk.Frame(parent_frame)
        accent_color_label = ttk.Label(accent_color_frame, text=self.lang.get_text('accent_color'))
        accent_color_label.pack(anchor=tk.W, pady=(10, 0))
        accent_color_button = ttk.Button(
            accent_color_frame, 
            text="Color: " + (self.settings.get('accent_color') or '#1f497d'),
            width=15,
            command=lambda: self.choose_color('accent_color', accent_color_button)
        )
        accent_color_button.pack(anchor=tk.W, pady=(5, 0))
        
        self.dynamic_widgets['ppt_accent_color'] = {
            'frame': accent_color_frame,
            'widgets': [accent_color_label, accent_color_button],
            'color_button': accent_color_button
        }
    
    def hex_to_rgb(self, hex_color):
        """Convert hex color to RGB tuple"""
        try:
            # Remove # if present
            hex_color = hex_color.lstrip('#')
            # Convert to RGB
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        except:
            # Return black as fallback
            return (0, 0, 0)
    
    def get_color_setting(self, setting_key, default_hex='#000000'):
        """Get color setting and return as RGB tuple"""
        hex_color = self.settings.get(setting_key, default_hex)
        return self.hex_to_rgb(hex_color)
    
    def choose_color(self, setting_key, button):
        """Open color chooser dialog and update the setting"""
        current_color = self.settings.get(setting_key, '#000000')
        color = colorchooser.askcolor(
            color=current_color,
            title=self.lang.get_text('choose_color')        )
        
        if color[1]:  # If a color was selected (not cancelled)
            # Update the setting
            self.settings.set(setting_key, color[1])
            # Update the button display with color preview
            self.update_color_button_display(button, color[1])
    def update_color_button_display(self, button, hex_color):
        """Update color button display with color preview"""
        # Show color code and visual indicator
        button.configure(text=f"Color: {hex_color}")
        
    def initialize_color_buttons(self):
        """Initialize color button displays with current settings"""
        if 'ppt_title_color' in self.dynamic_widgets:
            button = self.dynamic_widgets['ppt_title_color']['color_button']
            color = self.settings.get('title_text_color', '#000000')
            self.update_color_button_display(button, color)
            
        if 'ppt_bg_color' in self.dynamic_widgets:
            button = self.dynamic_widgets['ppt_bg_color']['color_button']
            color = self.settings.get('background_color', '#ffffff')
            self.update_color_button_display(button, color)
            
        if 'ppt_accent_color' in self.dynamic_widgets:
            button = self.dynamic_widgets['ppt_accent_color']['color_button']
            color = self.settings.get('accent_color', '#1f497d')
            self.update_color_button_display(button, color)
    
    def create_splitter_fields(self, parent_frame):
        """Create fields specific to image splitting"""
        
        # Rows per page
        rows_frame = ttk.Frame(parent_frame)
        rows_label = ttk.Label(rows_frame, text=self.lang.get_text('rows_per_page'))
        rows_label.pack(anchor=tk.W, pady=(10, 0))
        
        self.rows_per_page_var = tk.StringVar(value="10")
        rows_entry = ttk.Entry(rows_frame, textvariable=self.rows_per_page_var, width=10)
        rows_entry.pack(anchor=tk.W, pady=5)
        
        self.dynamic_widgets['splitter_rows'] = {
            'frame': rows_frame,
            'widgets': [rows_label, rows_entry]
        }
        
        # Maximum splits configuration
        splits_frame = ttk.Frame(parent_frame)
        splits_label = ttk.Label(splits_frame, text=self.lang.get_text('max_splits'))
        splits_label.pack(anchor=tk.W, pady=(10, 0))
        
        self.max_splits_var = tk.StringVar(value="2")
        splits_entry = ttk.Entry(splits_frame, textvariable=self.max_splits_var, width=10)
        splits_entry.pack(anchor=tk.W, pady=5)
        
        splits_help = ttk.Label(splits_frame, text=self.lang.get_text('max_splits_help'), 
                               font=("Arial", 8), foreground="gray")
        splits_help.pack(anchor=tk.W, pady=2)
        
        self.dynamic_widgets['splitter_max_splits'] = {
            'frame': splits_frame,
            'widgets': [splits_label, splits_entry, splits_help]
        }
        
        # Question detection
        question_frame = ttk.Frame(parent_frame)
        self.detect_questions_var = tk.BooleanVar(value=True)
        question_checkbox = ttk.Checkbutton(
            question_frame, 
            text=self.lang.get_text('detect_questions'),
            variable=self.detect_questions_var
        )
        question_checkbox.pack(anchor=tk.W, pady=(10, 0))
        
        self.dynamic_widgets['splitter_questions'] = {
            'frame': question_frame,
            'widgets': [question_checkbox]
        }
    
    def create_common_fields(self, parent_frame):
        """Create fields that might be used by multiple operations"""
        # This can be extended for fields used by multiple operations
        pass
    
    def on_operation_change(self, *args):
        """Handle operation selection changes to show/hide relevant fields"""
        operation = self.operation_var.get()
        
        # Hide all dynamic fields first
        for widget_group in self.dynamic_widgets.values():
            widget_group['frame'].pack_forget()
        
        # Show relevant fields based on operation
        if operation == "doc_to_pptx":
            # Show PowerPoint conversion fields
            self.dynamic_widgets['ppt_lines']['frame'].pack(fill=tk.X, pady=5)
            self.dynamic_widgets['ppt_spacing']['frame'].pack(fill=tk.X, pady=5)
            # Show color selection fields
            self.dynamic_widgets['ppt_title_color']['frame'].pack(fill=tk.X, pady=5)
            self.dynamic_widgets['ppt_bg_color']['frame'].pack(fill=tk.X, pady=5)
            self.dynamic_widgets['ppt_accent_color']['frame'].pack(fill=tk.X, pady=5)
            
        elif operation == "splitter":
            # Show splitter fields
            self.dynamic_widgets['splitter_rows']['frame'].pack(fill=tk.X, pady=5)
            self.dynamic_widgets['splitter_max_splits']['frame'].pack(fill=tk.X, pady=5)
            self.dynamic_widgets['splitter_questions']['frame'].pack(fill=tk.X, pady=5)
            
        # Other operations (title_cropper, blank_remover) don't need extra fields currently
        
        # Update UI text if language changed
        if hasattr(self, 'dynamic_widgets'):
            self.update_dynamic_field_text()
    
    def update_dynamic_field_text(self):
        """Update text for dynamic fields when language changes"""
        try:
            # Update PowerPoint conversion field labels
            if 'ppt_lines' in self.dynamic_widgets:
                widgets = self.dynamic_widgets['ppt_lines']['widgets']
                if len(widgets) >= 3:
                    widgets[0].config(text=self.lang.get_text('lines_per_slide'))
                    widgets[2].config(text=self.lang.get_text('lines_recommendation'))
            
            if 'ppt_spacing' in self.dynamic_widgets:
                widgets = self.dynamic_widgets['ppt_spacing']['widgets'] 
                if len(widgets) >= 3:
                    widgets[0].config(text=self.lang.get_text('line_spacing'))
                    widgets[2].config(text=self.lang.get_text('line_spacing_recommendation'))
            
            # Update splitter field labels
            if 'splitter_rows' in self.dynamic_widgets:
                widgets = self.dynamic_widgets['splitter_rows']['widgets']
                if len(widgets) >= 1:
                    widgets[0].config(text=self.lang.get_text('rows_per_page'))
            if 'splitter_questions' in self.dynamic_widgets:
                widgets = self.dynamic_widgets['splitter_questions']['widgets']
                if len(widgets) >= 1:
                    widgets[0].config(text=self.lang.get_text('detect_questions'))
            
            # Update color field labels
            if 'ppt_title_color' in self.dynamic_widgets:
                widgets = self.dynamic_widgets['ppt_title_color']['widgets']
                if len(widgets) >= 1:
                    widgets[0].config(text=self.lang.get_text('title_text_color'))
            
            if 'ppt_bg_color' in self.dynamic_widgets:
                widgets = self.dynamic_widgets['ppt_bg_color']['widgets']
                if len(widgets) >= 1:
                    widgets[0].config(text=self.lang.get_text('background_color'))
            
            if 'ppt_accent_color' in self.dynamic_widgets:
                widgets = self.dynamic_widgets['ppt_accent_color']['widgets']
                if len(widgets) >= 1:
                    widgets[0].config(text=self.lang.get_text('accent_color'))
            
        except Exception as e:
            print(f"Error updating dynamic field text: {e}")
    
    def find_optimal_split_points(self, img):
        """Find optimal split points that avoid cutting through text lines"""
        try:
            img_height, width = img.shape[:2]
            
            # Use OCR to detect text lines and their positions
            custom_config = r'-l tur --oem 3 --psm 6'
            data = pytesseract.image_to_data(img, config=custom_config, output_type=pytesseract.Output.DICT)
            
            # Extract text line boundaries
            text_lines = []
            for i, text in enumerate(data['text']):
                if text.strip():  # Non-empty text
                    top = data['top'][i]
                    height_text = data['height'][i]
                    bottom = top + height_text
                    text_lines.append((top, bottom))
            
            if not text_lines:
                # No text detected, fallback to simple middle split
                print("⚠️ No text detected, using fallback middle split")
                mid_height = img_height // 2
                return [(0, mid_height), (mid_height, img_height)]
            
            # Sort text lines by vertical position
            text_lines.sort()
            
            # Find gaps between text lines
            gaps = []
            for i in range(len(text_lines) - 1):
                gap_start = text_lines[i][1]  # Bottom of current line
                gap_end = text_lines[i + 1][0]  # Top of next line
                gap_size = gap_end - gap_start
                
                if gap_size > 20:  # Minimum gap size to consider for splitting
                    gap_center = (gap_start + gap_end) // 2
                    gaps.append((gap_center, gap_size))
            
            if not gaps:
                # No suitable gaps found, use fallback
                print("⚠️ No suitable text gaps found, using fallback middle split")
                mid_height = img_height // 2
                return [(0, mid_height), (mid_height, img_height)]
            
            # Find the gap closest to the middle of the image
            target_middle = img_height // 2
            best_gap = min(gaps, key=lambda gap: abs(gap[0] - target_middle))
            optimal_split_y = best_gap[0]
            
            print(f"📏 Found optimal split point at y={optimal_split_y} (target was y={target_middle})")
            print(f"✂️ Splitting image: [0:{optimal_split_y}] and [{optimal_split_y}:{img_height}]")
            
            return [(0, optimal_split_y), (optimal_split_y, img_height)]
            
        except Exception as e:
            print(f"⚠️ Error in intelligent split analysis: {e}")
            # Fallback to simple middle split
            img_height, width = img.shape[:2]
            mid_height = img_height // 2
            return [(0, mid_height), (mid_height, img_height)]
    
    def find_text_free_zones(self, img):
        """Find horizontal zones with minimal text content using both OCR and visual analysis"""
        try:
            img_height, width = img.shape[:2]
            
            # Method 1: OCR-based text detection
            custom_config = r'-l tur --oem 3 --psm 6'
            data = pytesseract.image_to_data(img, config=custom_config, output_type=pytesseract.Output.DICT)
            
            # Create a map of text density per row
            text_density = np.zeros(img_height)
            
            for i, text in enumerate(data['text']):
                if text.strip():
                    top = data['top'][i]
                    height_text = data['height'][i]
                    bottom = min(top + height_text, img_height - 1)
                    
                    # Add weight to rows containing text
                    for y in range(max(0, top), bottom + 1):
                        if y < img_height:
                            text_density[y] += len(text.strip())
            
            # Method 2: Visual analysis - detect horizontal lines with low content
            gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
            
            # Calculate row-wise variance (low variance = likely empty/uniform areas)
            row_variance = np.var(gray, axis=1)
            
            # Combine OCR text density with visual analysis
            combined_score = text_density + (row_variance / np.max(row_variance) * 50)
            
            # Smooth the score to avoid noise
            kernel_size = min(21, img_height // 10)  # Adaptive kernel size
            if kernel_size % 2 == 0:
                kernel_size += 1  # Ensure odd kernel size
            smoothed_score = cv2.GaussianBlur(combined_score.reshape(-1, 1), (1, kernel_size), 0).flatten()
            
            return smoothed_score
            
        except Exception as e:
            print(f"⚠️ Error in text-free zone analysis: {e}")
            return np.ones(img.shape[0])  # Return uniform scores as fallback
    
    def find_multiple_split_points(self, img, max_splits=3):
        """Find multiple optimal split points for complex images"""
        try:
            img_height, width = img.shape[:2]
            
            # Get text-free zones analysis
            content_scores = self.find_text_free_zones(img)
            
            # Find local minima (areas with least content)
            min_gap_size = img_height // 20  # Minimum 5% of image height
            potential_splits = []
            
            for i in range(min_gap_size, img_height - min_gap_size):
                # Check if this is a local minimum
                window_start = max(0, i - min_gap_size // 2)
                window_end = min(img_height, i + min_gap_size // 2)
                
                if content_scores[i] == np.min(content_scores[window_start:window_end]):
                    potential_splits.append((i, content_scores[i]))
            
            if not potential_splits:
                # No good splits found, fallback to simple middle split
                mid_height = img_height // 2
                return [(0, mid_height), (mid_height, img_height)]
            
            # Sort by content score (ascending - prefer areas with less content)
            potential_splits.sort(key=lambda x: x[1])
            
            # Select the best splits, ensuring they're well-spaced
            selected_splits = []
            min_distance = img_height // 4  # Minimum distance between splits
            
            for split_y, score in potential_splits:
                # Check if this split is far enough from already selected splits
                if all(abs(split_y - selected) >= min_distance for selected in selected_splits):
                    selected_splits.append(split_y)
                    if len(selected_splits) >= max_splits - 1:
                        break
            
            # Always include image boundaries and sort
            selected_splits.extend([0, img_height])
            selected_splits = sorted(set(selected_splits))
            
            # Create split ranges
            split_ranges = []
            for i in range(len(selected_splits) - 1):
                split_ranges.append((selected_splits[i], selected_splits[i + 1]))
            
            print(f"📐 Selected {len(split_ranges)} split ranges: {split_ranges}")
            
            return split_ranges
            
        except Exception as e:
            print(f"⚠️ Error in multiple split analysis: {e}")
            # Fallback to simple split
            mid_height = img_height // 2
            return [(0, mid_height), (mid_height, img_height)]
    
if __name__ == "__main__":
    root = tk.Tk()
    app = ImageCropperApp(root)
    root.mainloop()
